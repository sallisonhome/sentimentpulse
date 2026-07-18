"""Acceptance test: render_full_pack produces a 6-8 slide PPTX in both themes (v7.0).

Updated for the GTM Studio revisions (Median Commercial Potential, Commercial
Risks, Description & Razors added). v6.0 (2026-07-15) dropped the roadmap so
the pack lost the 6-slide roadmap block. v7.0 (2026-07-18 polish pass) made
slide count dynamic: the USP and Commercial Risks sub-decks each split into
2 slides when they carry 4-5 items (locked split rule: <=3 items = 1 slide,
4-5 items = [3, remainder] = 2 slides), so the assembled pack can now be
6, 7, or 8 slides depending on USP/risk counts.
"""

from __future__ import annotations

import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from gtm_pack import render_full_pack, render_pack_with_artifacts  # noqa: E402
from gtm_pack.sample_inputs import SAMPLE_INPUTS  # noqa: E402


@pytest.fixture(scope="module")
def out_dir(tmp_path_factory):
    return tmp_path_factory.mktemp("gtm_pack_out")


@pytest.mark.parametrize("theme", ["dark", "light"])
def test_render_full_pack_produces_expected_slide_count(theme, out_dir):
    """Slide count is dynamic (6-8) since v7.0's USP/Risks split rule.
    SAMPLE_INPUTS currently carries 4 USPs + 4 risks, which each split into
    2 slides, so this currently renders 8 slides -- but we assert a range
    so this test stays valid if SAMPLE_INPUTS' USP/risk counts change."""
    pptx = render_full_pack(SAMPLE_INPUTS, theme, out_dir)
    assert pptx.exists()
    prs = Presentation(str(pptx))
    assert 6 <= len(prs.slides) <= 8, (
        f"Expected 6-8 slides for {theme}, got {len(prs.slides)}"
    )


@pytest.mark.parametrize("theme", ["dark", "light"])
def test_no_residual_steam_refs(theme, out_dir):
    pptx = render_full_pack(SAMPLE_INPUTS, theme, out_dir)
    proc = subprocess.run(
        [sys.executable, "-m", "markitdown", str(pptx)],
        capture_output=True, text=True, timeout=60,
    )
    text = proc.stdout.lower()
    for forbidden in ("steam", "valve", "mad octopus", "steamworks", "aetheria"):
        assert forbidden not in text, f"Found forbidden '{forbidden}' in {theme} pack"


def test_render_pack_with_artifacts_produces_pdf_and_pngs(out_dir):
    """Full integration: PPTX + PDF + 6-8 PNGs (v7.0, dynamic slide count)."""
    result = render_pack_with_artifacts(SAMPLE_INPUTS, "dark", out_dir)
    assert result["pptx"].exists()
    assert result["pdf"].exists()
    assert 6 <= len(result["pngs"]) <= 8, (
        f"Expected 6-8 PNGs, got {len(result['pngs'])}"
    )
    for png in result["pngs"]:
        assert png.exists()
        assert png.stat().st_size > 10_000  # at least 10KB each
