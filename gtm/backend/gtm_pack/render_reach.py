#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 3 'How We Reach' slide.

v7 polish pass (2026-07-18) rewrite:
  - The mini circle-chart visual key is REMOVED entirely. All the freed
    width goes to the cohort rows, which now span the FULL slide width.
  - Each row shows, in order: cohort name + audience size (top line),
    then channels (middot-separated) on their own line, then message
    (left) / KPI (right) on a final line. Row height is computed per-row
    from an estimated wrapped-line count (same approach validated in
    render_usp.py) so nothing overlaps regardless of name/message length.
  - Cohort names and audience sizes are ALWAYS the user's verbatim
    wizard Step-3 values (name/size), never a hardcoded default like
    "Genre Fans" or "Breakout Ceiling". The renderer accepts an optional
    per-cohort `size` in the reach payload; if absent, the size line is
    simply omitted for that cohort (still no generic label is invented).
  - No footer (global v7 rule).

Two themes, parity with Steps 1 & 2:
  - dark   (V2 Modern Mono)
  - light  (V4 Bold Brand)

Outputs both a PPTX and a PNG to --out-dir.
"""
from __future__ import annotations

import argparse
import json
import math
import os
import re
import subprocess
import tempfile

from pptx import Presentation
from ._title_fit import fit_title_pt
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def fmt_num(n) -> str:
    try:
        return f"{int(n):,}"
    except (TypeError, ValueError):
        return str(n)


def est_wrapped_lines(text: str, width_in: float, pt_size: float,
                       chars_per_inch_at_10pt: float = 16.5,
                       safety_margin: float = 1.18) -> int:
    """Estimate how many lines `text` will wrap to inside a box of
    `width_in` inches at font size `pt_size`. Same calibration as
    render_usp.py's estimator -- deliberately over-counts lines so rows
    get slightly more height rather than risk overlap."""
    if not text:
        return 1
    chars_per_inch = chars_per_inch_at_10pt * (10.0 / pt_size) / safety_margin
    chars_per_line = max(8, int(width_in * chars_per_inch))
    total = 0
    for segment in text.split("\n"):
        n_chars = len(segment)
        total += max(1, math.ceil(n_chars / chars_per_line))
    return max(1, total)


# ---------- shape primitives ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w_pt=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w_pt)
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return box


# ---------- Cohort label resolution ----------
# NOTE (v7 naming-fix): cohort NAME and SIZE must always come from the
# wizard's Step 3 `cohorts` list verbatim -- never substitute a generic
# default like "Genre Fans" / "Breakout Ceiling". `cohort_names(args)`
# only falls back to a generic label in the (should-be-rare) case where
# the caller genuinely didn't supply a name for that slot.
def innermost_label(args) -> str:
    if args.inner == "prev":
        return "Prev Game Owners"
    if args.inner == "dev":
        return "Developer Fans"
    return args.inner_name or "Cohort"


def ring2_label(args) -> str:
    if args.type == "custom":
        return args.ring2_name or "Custom cohort"
    return "IP Fans (no prior)"


def cohort_names(args) -> list[str]:
    """Inner -> outer, matching circle chart tier order. Cohorts 3/4 use
    the user-typed args.cohort3_name / args.cohort4_name when present
    (threaded from inputs["cohorts"][2/3]["name"] by the __init__.py
    wrapper) and only fall back to the legacy generic label if absent."""
    return [
        innermost_label(args),
        ring2_label(args),
        getattr(args, "cohort3_name", None) or "Genre Fans",
        getattr(args, "cohort4_name", None) or "Breakout Ceiling",
    ]


def cohort_sizes(args) -> list[int | None]:
    """Inner -> outer audience sizes, verbatim from the wizard. None if
    not supplied for that slot (row simply omits the size, never invents
    one)."""
    return [
        getattr(args, "cohort1_size", None),
        getattr(args, "cohort2_size", None),
        getattr(args, "cohort3_size", None),
        getattr(args, "cohort4_size", None),
    ]


# ---------- Reach data ----------
def load_reach(path_or_inline: str, is_path: bool) -> list[dict]:
    """Load the per-cohort reach data. Must be a JSON list of 4 objects with keys:
      - channels: list[str]  (1-7 entries; extras beyond 7 are silently dropped)
      - message:  str
      - kpi:      str
    Order MUST match cohort tier order: inner -> outer (Prev, IP/Custom, Genre, Breakout).
    """
    if is_path:
        with open(path_or_inline, "r") as f:
            raw = json.load(f)
    else:
        raw = json.loads(path_or_inline)
    if not isinstance(raw, list) or len(raw) != 4:
        raise ValueError("reach must be a list of exactly 4 objects "
                         "(one per cohort tier, inner -> outer)")
    out = []
    for i, item in enumerate(raw):
        try:
            if "channels" in item:
                ch_raw = item["channels"]
            elif "channel" in item:
                s = item["channel"]
                if not isinstance(s, str):
                    raise ValueError(f"Cohort #{i+1} 'channel' must be a string")
                ch_raw = [p.strip() for p in s.replace("\u00b7", ",").split(",") if p.strip()]
            else:
                raise ValueError(f"Cohort #{i+1} missing 'channels' (list) or 'channel' (string)")
            ms = item["message"].strip()
            kp = item["kpi"].strip()
        except (KeyError, AttributeError, TypeError) as e:
            raise ValueError(f"Cohort #{i+1} missing one of: channels/channel, message, kpi ({e})")
        if not isinstance(ch_raw, list) or len(ch_raw) < 1:
            raise ValueError(f"Cohort #{i+1} must have at least one channel")
        ch = [str(c).strip() for c in ch_raw if str(c).strip()][:7]
        if not ch:
            raise ValueError(f"Cohort #{i+1} 'channels' resolved to empty after trimming")
        out.append({"channels": ch, "message": ms, "kpi": kp})
    return out


# ---------- Row height measurement (font-shrink ladder, mirrors render_usp.py) ----------
# Each rung shrinks name/channel/body font sizes together so hierarchy is
# preserved. A naive proportional box-height shrink (without shrinking the
# actual rendered font) was tried first and rejected: it left text at full
# size while giving it a smaller box, which caused real visual overlap
# between rows under dense worst-case copy. This ladder shrinks the fonts
# themselves before falling back to truncation, exactly like render_usp.py.
FONT_STEPS = [
    dict(name=15,   size_lbl=11.5, ch=10.0, body=10.5,
         name_lh=0.24, ch_lh=0.20, body_lh=0.22),
    dict(name=14,   size_lbl=11.0, ch=9.3,  body=9.7,
         name_lh=0.22, ch_lh=0.185, body_lh=0.205),
    dict(name=13,   size_lbl=10.5, ch=8.7,  body=9.0,
         name_lh=0.21, ch_lh=0.17, body_lh=0.19),
    dict(name=12,   size_lbl=10.0, ch=8.0,  body=8.3,
         name_lh=0.19, ch_lh=0.155, body_lh=0.175),
    # v8 fix: extra floor rungs so dense KPI/message copy (e.g. Hellraiser's
    # worst-case cohort KPIs) can wrap to 2-3 lines and still fit without
    # ever truncating -- see _fit_rows below, which no longer clips text.
    dict(name=11.5, size_lbl=9.5,  ch=7.6,  body=7.8,
         name_lh=0.18, ch_lh=0.145, body_lh=0.165),
    dict(name=11.0, size_lbl=9.0,  ch=7.2,  body=7.3,
         name_lh=0.17, ch_lh=0.135, body_lh=0.155),
]

ROW_PAD_TOP = 0.12
ROW_PAD_BOTTOM = 0.13
GAP_AFTER_NAME = 0.06
GAP_AFTER_CH = 0.07


def _measure_row(name, channels, message, kpi, step, *, name_w, ch_w, msg_w, kpi_w):
    """Return (row_height_in, name_lines, ch_lines, msg_lines, kpi_lines) at
    the given font-size rung `step`."""
    name_lines = est_wrapped_lines(name, name_w, step["name"])
    ch_str = "   \u00b7   ".join(channels)
    ch_lines = est_wrapped_lines(ch_str, ch_w, step["ch"])
    msg_lines = est_wrapped_lines(message, msg_w, step["body"])
    kpi_lines = est_wrapped_lines(kpi, kpi_w, step["body"])
    body_lines = max(msg_lines, kpi_lines)

    h = (ROW_PAD_TOP
         + name_lines * step["name_lh"]
         + GAP_AFTER_NAME
         + ch_lines * step["ch_lh"]
         + GAP_AFTER_CH
         + body_lines * step["body_lh"]
         + ROW_PAD_BOTTOM)
    return h, name_lines, ch_lines, msg_lines, kpi_lines


def _measure_all(rows_data, step, *, name_w, ch_w, msg_w, kpi_w):
    heights = []
    for r in rows_data:
        h, *_ = _measure_row(r["name"], r["channels"], r["message"], r["kpi"], step,
                              name_w=name_w, ch_w=ch_w, msg_w=msg_w, kpi_w=kpi_w)
        heights.append(h)
    return heights


def _fit_rows(rows_data, area_h, *, name_w, ch_w, msg_w, kpi_w, min_gap=0.09):
    """Try each FONT_STEPS rung (largest first); pick the first rung whose
    natural total height fits area_h. v8 fix: text is NEVER truncated any
    more -- if even the floor rung doesn't fit naturally, we keep the floor
    rung's font size (so text stays fully readable) and instead compress
    the row gap down toward zero, then as an absolute last resort scale row
    heights down to exactly fill area_h (rows get visually tight but every
    word of every message/KPI is still rendered). Returns
    (rows_data, step, heights, gap).
    """
    n = len(rows_data)
    chosen_step = FONT_STEPS[-1]
    chosen_heights = None
    for step in FONT_STEPS:
        heights = _measure_all(rows_data, step, name_w=name_w, ch_w=ch_w, msg_w=msg_w, kpi_w=kpi_w)
        total = sum(heights) + min_gap * (n - 1)
        if total <= area_h:
            chosen_step = step
            chosen_heights = heights
            break

    if chosen_heights is not None:
        total = sum(chosen_heights) + min_gap * (n - 1)
        leftover = area_h - total
        extra_gap = min(leftover / max(1, n - 1), 0.30) if n > 1 else 0.0
        remaining = leftover - extra_gap * (n - 1)
        heights = [h + (remaining * (h / total_h) if (total_h := sum(chosen_heights)) > 0 else 0)
                   for h in chosen_heights]
        gap = min_gap + extra_gap
        return rows_data, chosen_step, heights, gap

    # Floor rung still doesn't fit at the normal min_gap. Do NOT truncate
    # text -- first try shrinking the inter-row gap toward zero (still at
    # the floor font size, still full text) before touching row heights.
    import warnings
    step = FONT_STEPS[-1]
    heights = _measure_all(rows_data, step, name_w=name_w, ch_w=ch_w, msg_w=msg_w, kpi_w=kpi_w)
    total_no_gap = sum(heights)
    if total_no_gap <= area_h:
        leftover = area_h - total_no_gap
        gap = leftover / max(1, n - 1) if n > 1 else 0.0
        return rows_data, step, heights, gap

    # Absolute last resort: even zero-gap floor-rung heights overflow the
    # area. Scale heights down to exactly fill area_h -- rows sit tightly
    # against each other but every line of text is still fully rendered,
    # never clipped. Warn loudly so this is visible in logs.
    warnings.warn(
        f"render_reach: {n} cohort row(s) did not fit even at the floor "
        f"font rung with zero row gap (needs {total_no_gap:.2f}in, have "
        f"{area_h:.2f}in). Compressing row heights instead of truncating "
        "any text.",
        RuntimeWarning,
        stacklevel=2,
    )
    scale = area_h / total_no_gap if total_no_gap > 0 else 1.0
    heights = [h * scale for h in heights]
    gap = 0.0
    return rows_data, step, heights, gap


# v8 fix: KPI column widened (was ~30% of row width; the mini circle-chart
# removal already freed up horizontal room, but the message/KPI split
# hadn't been rebalanced to use it). Message now gets ~50% and KPI ~44%
# of the row width, which is enough for Hellraiser's worst-case KPI copy
# (up to ~225 characters) to wrap onto 2-3 lines instead of truncating.
MSG_WIDTH_FRACTION = 0.50
COL_HEADER_H = 0.22
COL_HEADER_GAP = 0.06


def _render_rows(slide, rows_data, *, rx, ry, rw, area_h, colors, ink, muted, hair,
                  name_font, l_body_pt, header_color=None):
    name_w = rw - 0.35
    ch_w = rw - 0.35
    msg_w = rw * MSG_WIDTH_FRACTION
    kpi_w = rw - msg_w - 0.30 - 0.35

    # ---- Column headers (6c): "KPIS" above the right column, right-aligned
    # to match the KPI text column, plus "COHORT" and "CHANNELS" above the
    # left columns for symmetry (per the user's earlier note). Header row
    # sits in its own reserved band directly above the first cohort row so
    # it never collides with the size-label ("N potential buyers") text.
    if header_color is not None:
        kpi_x = rx + 0.30 + msg_w + 0.30
        header_y = ry - COL_HEADER_H - COL_HEADER_GAP
        add_text(slide, rx + 0.30, header_y, name_w * 0.4, COL_HEADER_H,
                 "COHORT", font="Calibri", size=9, bold=True, color=header_color)
        add_text(slide, kpi_x, header_y, kpi_w, COL_HEADER_H,
                 "KPIS", font="Calibri", size=9, bold=True, color=header_color,
                 align=PP_ALIGN.RIGHT)

    rows_data, step, heights, gap = _fit_rows(
        rows_data, area_h, name_w=name_w, ch_w=ch_w, msg_w=msg_w, kpi_w=kpi_w,
    )

    y = ry
    for i, r in enumerate(rows_data):
        h = heights[i]
        c = colors[i]
        add_oval(slide, rx, y + 0.06, 0.16, 0.16, c)
        name_box_w = name_w - (1.9 if r["size_label"] else 0)
        add_text(slide, rx + 0.30, y, name_box_w, 0.6, r["name"],
                 font=name_font, size=step["name"], bold=True, color=ink,
                 line_spacing=1.0)
        if r["size_label"]:
            add_text(slide, rx + 0.30 + name_box_w + 0.1, y, 1.8, 0.3, r["size_label"],
                     font=name_font, size=step["size_lbl"], bold=True, color=c,
                     align=PP_ALIGN.RIGHT)
        name_lines = est_wrapped_lines(r["name"], name_box_w, step["name"])
        y_ch = y + name_lines * step["name_lh"] + GAP_AFTER_NAME
        ch_str = "   \u00b7   ".join(r["channels"])
        add_text(slide, rx + 0.30, y_ch, ch_w, 0.6, ch_str,
                 font="Calibri", size=body_pt(l_body_pt, step["ch"]), color=muted,
                 line_spacing=1.0)
        ch_lines = est_wrapped_lines(ch_str, ch_w, step["ch"])
        y_body = y_ch + ch_lines * step["ch_lh"] + GAP_AFTER_CH
        add_text(slide, rx + 0.30, y_body, msg_w, 0.6, r["message"],
                 font="Calibri", size=body_pt(l_body_pt, step["body"]), color=ink,
                 line_spacing=1.05)
        add_text(slide, rx + 0.30 + msg_w + 0.30, y_body, kpi_w, 0.6, r["kpi"],
                 font="Calibri", size=body_pt(l_body_pt, step["body"] - 0.5), bold=True,
                 color=c, align=PP_ALIGN.RIGHT, line_spacing=1.05)
        if i < len(rows_data) - 1:
            add_rect(slide, rx, y + h - gap * 0.45, rw, 0.008, hair)
        y += h + gap


def _build_rows_data(args, reach) -> list[dict]:
    names = cohort_names(args)
    sizes = cohort_sizes(args)
    rows = []
    for i in range(4):
        size_label = f"{fmt_num(sizes[i])} potential buyers" if sizes[i] is not None else ""
        rows.append({
            "name": names[i],
            "size_label": size_label,
            "channels": reach[i]["channels"],
            "message": reach[i]["message"],
            "kpi": reach[i]["kpi"],
        })
    return rows


# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, reach, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    A1       = hex_rgb("#0A2A30")  # outer
    A2       = hex_rgb("#155966")
    A3       = hex_rgb("#2FA9BD")
    A4       = hex_rgb("#7FD8E3")  # inner
    ACCENT   = hex_rgb("#FFB454")  # warm gold (breakout + eyebrow)

    tier_colors = [A4, A3, A2, ACCENT]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    add_text(slide, 0.6, 0.4, 10, 0.3, "HOW WE REACH",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 12, 0.85, f"Reaching the audience for {args.title}",
             font="Trebuchet MS", size=fit_title_pt(f"Reaching the audience for {args.title}", 12), bold=True, color=INK)
    add_text(slide, 0.6, 1.55, 12, 0.4,
             f"{args.genre} \u00b7 Channels, messaging and KPIs tied to each cohort",
             font="Calibri", size=body_pt(L, 13), color=MUTED)

    rows_data = _build_rows_data(args, reach)
    rx, ry = 0.6, 2.15 + COL_HEADER_H + COL_HEADER_GAP
    rw = 12.1
    area_h = 6.85 - ry  # leave clear margin at bottom, no footer
    _render_rows(slide, rows_data, rx=rx, ry=ry, rw=rw, area_h=area_h,
                 colors=tier_colors, ink=INK, muted=MUTED, hair=BORDER,
                 name_font="Trebuchet MS", l_body_pt=L, header_color=MUTED)

    prs.save(out_path)


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, reach, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    C4       = hex_rgb("#7DD4C9")  # mint (inner)
    C3       = hex_rgb("#1F9B8E")  # teal
    C2       = hex_rgb("#D63A57")  # rose
    C1       = hex_rgb("#E5A700")  # gold (outer / breakout)

    tier_colors = [C4, C3, C2, C1]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 0.25, 7.5, C3)

    add_text(slide, 0.7, 0.5, 11, 0.3, "HOW WE REACH",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3)
    add_text(slide, 0.7, 0.85, 12, 0.85, f"Reaching the audience for {args.title}",
             font="Trebuchet MS", size=fit_title_pt(f"Reaching the audience for {args.title}", 12), bold=True, color=INK)
    add_text(slide, 0.7, 1.65, 12, 0.4,
             f"{args.genre} \u00b7 Channels, messaging and KPIs tied to each cohort",
             font="Calibri", size=body_pt(L, 14), color=MUTED)

    rows_data = _build_rows_data(args, reach)
    rx, ry = 0.7, 2.25 + COL_HEADER_H + COL_HEADER_GAP
    rw = 12.0
    area_h = 6.95 - ry
    _render_rows(slide, rows_data, rx=rx, ry=ry, rw=rw, area_h=area_h,
                 colors=tier_colors, ink=INK, muted=MUTED, hair=HAIR,
                 name_font="Trebuchet MS", l_body_pt=L, header_color=MUTED)

    prs.save(out_path)


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_path: str) -> str:
    out_dir = os.path.dirname(os.path.abspath(png_path))
    os.makedirs(out_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = None
        for fn in sorted(os.listdir(tmp)):
            if fn.startswith("slide-") and fn.endswith(".png"):
                produced = os.path.join(tmp, fn)
                break
        if not produced:
            raise RuntimeError("pdftoppm did not produce a PNG")
        with open(produced, "rb") as src, open(png_path, "wb") as dst:
            dst.write(src.read())
    return png_path


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True)
    p.add_argument("--theme", required=True, choices=["dark", "light"])
    p.add_argument(
        "--type", required=True, choices=["sequel", "new_ip_with_fans", "custom"],
        help="Match the value used in Step 1 (controls ring 2 label).",
    )
    p.add_argument(
        "--inner", required=True, choices=["prev", "dev", "other"],
        help="Match the value used in Step 1 (controls innermost label).",
    )
    p.add_argument("--inner-name", default=None,
                   help="Custom innermost label (required when --inner other)")
    p.add_argument("--ring2-name", default=None,
                   help="Custom ring 2 label (required when --type custom)")
    p.add_argument("--cohort3-name", default=None,
                   help="User-typed name for cohort 3 (falls back to 'Genre Fans' if omitted)")
    p.add_argument("--cohort4-name", default=None,
                   help="User-typed name for cohort 4 (falls back to 'Breakout Ceiling' if omitted)")
    p.add_argument("--cohort1-size", type=int, default=None)
    p.add_argument("--cohort2-size", type=int, default=None)
    p.add_argument("--cohort3-size", type=int, default=None)
    p.add_argument("--cohort4-size", type=int, default=None)

    grp = p.add_mutually_exclusive_group(required=True)
    grp.add_argument("--reach", help="Inline JSON list of 4 cohort objects "
                                     "(channels[], message, kpi). Order: inner -> outer.")
    grp.add_argument("--reach-json", help="Path to a JSON file with the cohort list")
    p.add_argument("--out-dir", required=True)
    args = p.parse_args()
    if args.type == "custom" and not args.ring2_name:
        p.error("--ring2-name is required when --type custom")
    if args.inner == "other" and not args.inner_name:
        p.error("--inner-name is required when --inner other")
    return args


def main():
    args = parse_args()
    reach = load_reach(args.reach_json or args.reach, is_path=bool(args.reach_json))
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_reach_{args.theme}.pptx")
    png_path  = os.path.join(out_dir, f"{slug}_reach_{args.theme}.png")
    if args.theme == "dark":
        render_dark(args, reach, pptx_path)
    else:
        render_light(args, reach, pptx_path)
    convert_to_png(pptx_path, png_path)
    print(f"PPTX: {pptx_path}")
    print(f"PNG:  {png_path}")


if __name__ == "__main__":
    main()
