#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 5 'Median Commercial Potential' slide.

Renders in output-position 2 (right after the Sizing Rings slide), even
though it is internally numbered Step 5.

Two themes, parity with Steps 1-4:
  - dark   (V2 Modern Mono)  : dark slide, teal ramp + warm gold accent
  - light  (V4 Bold Brand)   : light slide, mint/teal/rose/gold tier ramp, left teal stripe

Layout:
  - Left column: 4 KPI cards (median revenue, median units, avg price, avg hours)
  - Right column: platform revenue projection table (one row per selected platform)
  - Footer strip below tables with comp-set attribution note

Outputs both a PPTX and a PNG to --out-dir.
"""
from __future__ import annotations

import argparse
import os
import re
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


# ---------- shape primitives (shared pattern) ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return box


# ---------- platform projection math (LOCKED) ----------
BASE_WEIGHTS = {"PC": 50, "PS5": 30, "XSX": 12, "SWITCH2": 8}
PLATFORM_ORDER = ["PC", "PS5", "XSX", "SWITCH2"]


def compute_platform_shares(selected: list[str]) -> list[tuple[str, float]]:
    """Return [(platform, share_pct), ...] sorted by base weight descending,
    for platforms in `selected` only."""
    sel = [p for p in PLATFORM_ORDER if p in selected]
    if not sel:
        # ValueError (not SystemExit) so the FastAPI endpoint can catch
        # it as a data-validation error and return 400 instead of 500.
        raise ValueError("At least one platform must be selected")
    total_selected = sum(BASE_WEIGHTS[p] for p in sel)
    return [(p, BASE_WEIGHTS[p] / total_selected * 100.0) for p in sel]


# ---------- formatting ----------
# DISPLAY FORMAT (locked): the underlying value is ALWAYS millions of dollars
# (median_revenue_usd_millions). Both the KPI card and the projection table
# render the bare number as "$X.XX" (no "M" suffix) and pair it with a small
# muted "in millions" / "shown in millions of USD" subcopy label so the unit
# stays legible without cluttering every cell. Do not re-add an "M" suffix to
# the number itself -- the unit lives in the subcopy line only.
def fmt_dollars_2dp(value: float) -> str:
    """Format a plain float as $X.XX (two decimals). Used for both the KPI
    card's median revenue number and the per-platform projection table cells
    -- both values are in millions of dollars, with the unit communicated via
    a separate subcopy label rather than a suffix on the number."""
    return f"${value:.2f}"


# Backward-compat aliases (older skill/back-end call sites may still import
# these names). Both now delegate to the same plain-$X.XX formatter -- the
# "M" suffix formatting is retired per the display-format refinement.
def fmt_millions(millions: float) -> str:
    return fmt_dollars_2dp(millions)


def fmt_millions_precise(millions: float) -> str:
    return fmt_dollars_2dp(millions)


def fmt_units(units: int) -> str:
    """Format a raw unit count as X.XM units (>=1M), XXK units (>=1K), or raw count."""
    if units >= 1_000_000:
        return f"{units / 1_000_000.0:.1f}M units"
    if units >= 1_000:
        return f"{units / 1_000.0:.0f}K units"
    return f"{units} units"


def fmt_price(dollars: float) -> str:
    return f"${dollars:.2f}"


# ---------- data loading ----------
def parse_platforms(raw: str) -> list[str]:
    parts = [p.strip().upper() for p in raw.split(",") if p.strip()]
    valid = set(PLATFORM_ORDER)
    bad = [p for p in parts if p not in valid]
    if bad:
        raise ValueError(f"Unknown platform(s): {bad}. Valid: {PLATFORM_ORDER}")
    if not parts:
        raise ValueError("platforms must include at least one platform")
    # De-dup while preserving canonical order
    return [p for p in PLATFORM_ORDER if p in set(parts)]


# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    SURFACE  = hex_rgb("#161A21")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    A3       = hex_rgb("#2FA9BD")
    A4       = hex_rgb("#7FD8E3")
    ACCENT   = hex_rgb("#FFB454")  # warm gold

    platforms = parse_platforms(args.platforms)
    shares = compute_platform_shares(platforms)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    add_text(slide, 0.6, 0.4, 10, 0.3, "STEP 05 \u00b7 MEDIAN COMMERCIAL POTENTIAL",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 12, 0.85, args.title,
             font="Trebuchet MS", size=34, bold=True, color=INK)
    add_text(slide, 0.6, 1.55, 12, 0.4,
             f"Genre benchmark \u00b7 {args.genre} \u00b7 comp set: {args.comp_set_name}",
             font="Calibri", size=body_pt(L, 13), color=MUTED)

    # ---- Hero: median revenue is the headline number, set in plain
    # typographic hierarchy (no card/border) so it reads as THE answer to
    # "what can this game make", with the other three KPIs demoted to a
    # calmer supporting row underneath, separated only by hairlines.
    hero_x, hero_y = 0.6, 2.35
    add_text(slide, hero_x, hero_y, 6.0, 0.3, "MEDIAN REVENUE \u00b7 COMP SET",
             font="Calibri", size=body_pt(L, 10), bold=True, color=MUTED)
    add_text(slide, hero_x, hero_y + 0.32, 6.0, 0.95, fmt_dollars_2dp(args.median_revenue_usd_millions),
             font="Trebuchet MS", size=54, bold=True, color=A4)
    add_text(slide, hero_x, hero_y + 1.28, 6.0, 0.3, "in millions of USD",
             font="Calibri", size=body_pt(L, 10), color=MUTED)

    # Supporting KPI row: three simple stat columns, hairline above only
    sup_y = hero_y + 1.72
    add_rect(slide, hero_x, sup_y, 5.9, 0.008, BORDER)
    sup_kpis = [
        ("MEDIAN UNITS SOLD", fmt_units(args.median_units_sold)),
        ("AVG PRICE", fmt_price(args.avg_price_usd)),
        ("AVG HOURS PLAYED", f"{args.avg_hours_played:.1f} hrs"),
    ]
    sup_col_w = 1.97
    for i, (label, val) in enumerate(sup_kpis):
        x = hero_x + i * sup_col_w
        add_text(slide, x, sup_y + 0.2, sup_col_w - 0.2, 0.55, label,
                 font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
        add_text(slide, x, sup_y + 0.52, sup_col_w - 0.2, 0.45, val,
                 font="Trebuchet MS", size=19, bold=True, color=INK)

    # ---- Right column: platform projection table ----
    # Hairline-row table, no header box, no alternating shading -- reads
    # calmer and lets the revenue column (right-aligned) breathe.
    rx, ry = 7.05, 2.35
    rw = 5.65
    add_text(slide, rx, ry, 2.4, 0.3, "PROJECTED BY PLATFORM",
             font="Calibri", size=body_pt(L, 10), bold=True, color=MUTED)
    header_h = 0.5
    headers = ["PLATFORM", "SHARE", "REVENUE", "UNITS"]
    col_x = [rx, rx + 1.55, rx + 2.75, rx + 4.15]
    col_w = [1.45, 1.1, 1.3, 1.5]
    hy = ry + 0.4
    for h_text, cx, cw in zip(headers, col_x, col_w):
        add_text(slide, cx, hy, cw, 0.25, h_text,
                 font="Calibri", size=body_pt(L, 8), bold=True, color=ACCENT,
                 align=(PP_ALIGN.LEFT if h_text == "PLATFORM" else PP_ALIGN.RIGHT))
    add_rect(slide, rx, hy + 0.3, rw, 0.012, BORDER)

    row_h = 0.62
    table_top = hy + 0.42
    for i, (plat, share) in enumerate(shares):
        y = table_top + i * row_h
        rev_millions = args.median_revenue_usd_millions * (share / 100.0)
        units = int(round(args.median_units_sold * (share / 100.0)))
        add_text(slide, col_x[0], y + 0.1, col_w[0], 0.35, plat,
                 font="Trebuchet MS", size=15, bold=True, color=INK)
        add_text(slide, col_x[1], y + 0.14, col_w[1], 0.3, f"{share:.1f}%",
                 font="Calibri", size=body_pt(L, 11), color=MUTED, align=PP_ALIGN.RIGHT)
        add_text(slide, col_x[2], y + 0.1, col_w[2], 0.35, fmt_dollars_2dp(rev_millions),
                 font="Trebuchet MS", size=15, bold=True, color=A4, align=PP_ALIGN.RIGHT)
        add_text(slide, col_x[3], y + 0.14, col_w[3], 0.3, fmt_units(units),
                 font="Calibri", size=body_pt(L, 11), color=MUTED, align=PP_ALIGN.RIGHT)
        if i < len(shares) - 1:
            add_rect(slide, rx, y + row_h - 0.06, rw, 0.008, BORDER)

    table_bottom = table_top + len(shares) * row_h

    # Unit-clarity + attribution, combined into one calmer line
    add_text(slide, rx, table_bottom + 0.12, rw, 0.5,
             (f"Revenue in millions of USD. Anchored by Genre Pulse comp-set "
              f"medians ({args.comp_set_titles} titles); PC-only caps at 100% "
              f"of PC-only potential."),
             font="Calibri", size=body_pt(L, 8), italic=True, color=MUTED)

    # Standard footer
    add_text(slide, 0.6, 7.1, 12, 0.25,
             "GTM SLIDE PACK \u00b7 STEP 05",
             font="Calibri", size=body_pt(L, 8), bold=True, color=MUTED)

    prs.save(out_path)


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    C3       = hex_rgb("#1F9B8E")  # structural teal
    C4       = hex_rgb("#7DD4C9")  # mint
    C2       = hex_rgb("#D63A57")  # rose

    platforms = parse_platforms(args.platforms)
    shares = compute_platform_shares(platforms)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 0.25, 7.5, C3)

    add_text(slide, 0.7, 0.5, 11, 0.3, "STEP 05 \u00b7 MEDIAN COMMERCIAL POTENTIAL",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3)
    add_text(slide, 0.7, 0.85, 12, 0.85, args.title,
             font="Trebuchet MS", size=34, bold=True, color=INK)
    add_text(slide, 0.7, 1.65, 12, 0.4,
             f"Genre benchmark \u00b7 {args.genre} \u00b7 comp set: {args.comp_set_name}",
             font="Calibri", size=body_pt(L, 14), color=MUTED)

    # ---- Hero: median revenue is the headline number ----
    hero_x, hero_y = 0.7, 2.45
    add_text(slide, hero_x, hero_y, 6.0, 0.3, "MEDIAN REVENUE \u00b7 COMP SET",
             font="Calibri", size=body_pt(L, 10), bold=True, color=MUTED)
    add_text(slide, hero_x, hero_y + 0.32, 6.0, 0.95, fmt_dollars_2dp(args.median_revenue_usd_millions),
             font="Trebuchet MS", size=54, bold=True, color=C2)
    add_text(slide, hero_x, hero_y + 1.28, 6.0, 0.3, "in millions of USD",
             font="Calibri", size=body_pt(L, 10), color=MUTED)

    # Supporting KPI row: three simple stat columns, hairline above only
    sup_y = hero_y + 1.72
    add_rect(slide, hero_x, sup_y, 5.8, 0.008, HAIR)
    sup_kpis = [
        ("MEDIAN UNITS SOLD", fmt_units(args.median_units_sold)),
        ("AVG PRICE", fmt_price(args.avg_price_usd)),
        ("AVG HOURS PLAYED", f"{args.avg_hours_played:.1f} hrs"),
    ]
    sup_col_w = 1.93
    for i, (label, val) in enumerate(sup_kpis):
        x = hero_x + i * sup_col_w
        add_text(slide, x, sup_y + 0.2, sup_col_w - 0.2, 0.55, label,
                 font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
        add_text(slide, x, sup_y + 0.52, sup_col_w - 0.2, 0.45, val,
                 font="Trebuchet MS", size=19, bold=True, color=INK)

    # ---- Right column: platform projection table ----
    rx, ry = 7.1, 2.45
    rw = 5.5
    add_text(slide, rx, ry, 2.4, 0.3, "PROJECTED BY PLATFORM",
             font="Calibri", size=body_pt(L, 10), bold=True, color=MUTED)
    headers = ["PLATFORM", "SHARE", "REVENUE", "UNITS"]
    col_x = [rx, rx + 1.5, rx + 2.65, rx + 4.0]
    col_w = [1.4, 1.05, 1.25, 1.5]
    hy = ry + 0.4
    for h_text, cx, cw in zip(headers, col_x, col_w):
        add_text(slide, cx, hy, cw, 0.25, h_text,
                 font="Calibri", size=body_pt(L, 8), bold=True, color=C3,
                 align=(PP_ALIGN.LEFT if h_text == "PLATFORM" else PP_ALIGN.RIGHT))
    add_rect(slide, rx, hy + 0.3, rw, 0.012, HAIR)

    row_h = 0.62
    table_top = hy + 0.42
    for i, (plat, share) in enumerate(shares):
        y = table_top + i * row_h
        rev_millions = args.median_revenue_usd_millions * (share / 100.0)
        units = int(round(args.median_units_sold * (share / 100.0)))
        add_text(slide, col_x[0], y + 0.1, col_w[0], 0.35, plat,
                 font="Trebuchet MS", size=15, bold=True, color=INK)
        add_text(slide, col_x[1], y + 0.14, col_w[1], 0.3, f"{share:.1f}%",
                 font="Calibri", size=body_pt(L, 11), color=MUTED, align=PP_ALIGN.RIGHT)
        add_text(slide, col_x[2], y + 0.1, col_w[2], 0.35, fmt_dollars_2dp(rev_millions),
                 font="Trebuchet MS", size=15, bold=True, color=C2, align=PP_ALIGN.RIGHT)
        add_text(slide, col_x[3], y + 0.14, col_w[3], 0.3, fmt_units(units),
                 font="Calibri", size=body_pt(L, 11), color=MUTED, align=PP_ALIGN.RIGHT)
        if i < len(shares) - 1:
            add_rect(slide, rx, y + row_h - 0.06, rw, 0.008, HAIR)

    table_bottom = table_top + len(shares) * row_h

    add_text(slide, rx, table_bottom + 0.12, rw, 0.5,
             (f"Revenue in millions of USD. Anchored by Genre Pulse comp-set "
              f"medians ({args.comp_set_titles} titles); PC-only caps at 100% "
              f"of PC-only potential."),
             font="Calibri", size=body_pt(L, 8), italic=True, color=MUTED)

    add_text(slide, 0.7, 7.1, 12, 0.25,
             "GTM Slide Pack \u00b7 Step 05",
             font="Calibri", size=body_pt(L, 9), color=MUTED)

    prs.save(out_path)


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_path: str) -> str:
    out_dir = os.path.dirname(os.path.abspath(png_path))
    os.makedirs(out_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = None
        for fn in sorted(os.listdir(tmp)):
            if fn.startswith("slide-") and fn.endswith(".png"):
                produced = os.path.join(tmp, fn)
                break
        if not produced:
            raise RuntimeError("pdftoppm did not produce a PNG")
        with open(produced, "rb") as src, open(png_path, "wb") as dst:
            dst.write(src.read())
    return png_path


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--theme", required=True, choices=["dark", "light"])
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True, help="The game's genre (plain text)")
    p.add_argument("--comp-set-name", required=True,
                   help='Genre Pulse comp set label, e.g. "Horror \u2014 19 titles"')
    p.add_argument("--median-revenue-usd-millions", type=float, required=True,
                   dest="median_revenue_usd_millions",
                   help="Median revenue from Genre Pulse aggregate, in MILLIONS of dollars (e.g. 4.7 for $4.7M)")
    p.add_argument("--median-units-sold", type=int, required=True,
                   help="Median units sold from Genre Pulse aggregate (raw unit count)")
    p.add_argument("--avg-price-usd", type=float, required=True,
                   help="Average price from Genre Pulse aggregate, in dollars (decimal)")
    p.add_argument("--avg-hours-played", type=float, required=True,
                   help="Average hours played from Genre Pulse aggregate")
    p.add_argument("--platforms", required=True,
                   help='Comma-separated subset of PC,PS5,XSX,SWITCH2')
    p.add_argument("--out-dir", required=True)
    args = p.parse_args()

    # Derive title count for footer strip, e.g. "Horror \u2014 19 titles" -> 19
    m = re.search(r"(\d+)", args.comp_set_name)
    args.comp_set_titles = m.group(1) if m else "N"

    return args


def main():
    args = parse_args()
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_commercial_potential_{args.theme}.pptx")
    png_path  = os.path.join(out_dir, f"{slug}_commercial_potential_{args.theme}.png")

    if args.theme == "dark":
        render_dark(args, pptx_path)
    else:
        render_light(args, pptx_path)

    convert_to_png(pptx_path, png_path)
    print(f"PPTX: {pptx_path}")
    print(f"PNG:  {png_path}")


if __name__ == "__main__":
    main()
