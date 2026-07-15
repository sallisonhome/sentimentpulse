"""GTM Slide Pack rendering engine.

Lifted from the `gtm-slide-pack-kickoff` user skill (v5.0, currency-corrected
revision — see gtm_revisions_summary.md for full history).
Exposes callable wrapper functions over the original argparse-driven renderers.

The original render_*.py modules are CLI scripts that parse argparse.Namespace.
This module wraps each one with a function that accepts a dict of inputs,
constructs a Namespace, and invokes the underlying renderer functions directly
(bypassing argparse and sys.exit).

Full 12-slide pack order (output position, NOT internal step numbers):
  1. Sizing Circle           (render_sizing_circle)
  2. Median Commercial Potential (render_commercial_potential) -- NEW
  3. USP / Pillars           (render_usp)
  4. How We Reach            (render_reach)
  5-10. Roadmap 4.1-4.6      (render_roadmap)
  11. Commercial Risks       (render_commercial_risks) -- NEW
  12. Description & Razors   (render_description_razors) -- NEW

Currency units (do not confuse these -- see gtm_revisions_summary.md):
  - median_revenue_usd_millions: float, MILLIONS of dollars (e.g. 4.7 == $4.7M)
  - avg_price_usd: float, PLAIN dollars (e.g. 39.99)
  - median_units_sold: int, RAW unit count (e.g. 1782675)
"""

from __future__ import annotations

import datetime as dt
import json
import os
import re
import shutil
import subprocess
import tempfile
from argparse import Namespace
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

# Import the underlying renderer modules. The skill modules use absolute paths
# for asset loading — we patch that below.
from . import render_sizing_circle as _sizing
from . import render_usp as _usp
from . import render_reach as _reach
from . import render_roadmap as _roadmap
from . import render_commercial_potential as _commercial_potential
from . import render_commercial_risks as _commercial_risks
from . import render_description_razors as _description_razors

# Path to bundled assets (roadmap_phases.json etc.)
ASSETS_DIR = Path(__file__).resolve().parent / "assets"


# ── Helpers ───────────────────────────────────────────────────────────────────


def _slugify(value: str) -> str:
    return _sizing.slugify(value)


def _convert_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    """Convert PPTX to PDF using LibreOffice (soffice). Returns PDF path."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    pdf_name = pptx_path.stem + ".pdf"
    pdf_path = out_dir / pdf_name
    # soffice writes to --outdir using same stem
    subprocess.run(
        [
            "soffice", "--headless", "--convert-to", "pdf",
            "--outdir", str(out_dir), str(pptx_path),
        ],
        check=True,
        capture_output=True,
        timeout=120,
    )
    if not pdf_path.exists():
        raise RuntimeError(f"PDF conversion failed: {pdf_path} not produced")
    return pdf_path


def _pdf_to_pngs(pdf_path: Path, out_dir: Path, dpi: int = 110, prefix: str = "slide") -> list[Path]:
    """Render PDF pages to PNGs via pdftoppm. Returns list of PNG paths in order."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["pdftoppm", "-png", "-r", str(dpi), str(pdf_path), str(out_dir / prefix)],
        check=True,
        capture_output=True,
        timeout=120,
    )
    pngs = sorted(out_dir.glob(f"{prefix}-*.png"))
    return pngs


# ── Individual renderers (callable wrappers) ─────────────────────────────────


def render_sizing_circle(inputs: dict[str, Any], theme: str, out_dir: Path,
                          language: str = "en") -> Path:
    """Render Step 1 — Sizing Circle. Returns PPTX path."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    cohort_sizes = inputs["cohorts"]  # list of 4 dicts with size
    inner_choice = inputs.get("inner", "other")
    game_type = inputs.get("game_type", "custom")

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        type=game_type,
        inner=inner_choice,
        inner_name=inputs["cohorts"][0]["name"] if inner_choice == "other" else None,
        inner_definition=inputs.get("inner_definition"),
        prev=int(cohort_sizes[0]["size"]),
        ip_fans=int(cohort_sizes[1]["size"]),
        genre_fans=int(cohort_sizes[2]["size"]),
        breakout=int(cohort_sizes[3]["size"]),
        ring2_name=inputs["cohorts"][1]["name"] if game_type == "custom" else None,
        ring2_definition=inputs.get("ring2_definition"),
        theme=theme,
        out_dir=str(out_dir),
        language=language,
    )

    pptx_path = out_dir / f"{slug}_sizing_circle_{theme}.pptx"
    if theme == "dark":
        _sizing.render_dark(args, str(pptx_path))
    else:
        _sizing.render_light(args, str(pptx_path))
    return pptx_path


def render_usp(inputs: dict[str, Any], theme: str, out_dir: Path,
                language: str = "en") -> Path:
    """Render Step 2 — USP Manifesto. Returns PPTX path."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    # Build args namespace; load_usps reads from args.usps (JSON string) or args.usps_json (path)
    usps_payload = inputs["usps"]  # list of {title/description/support OR proof, strategy, enabled}
    # Normalize "support" -> "proof" for skill compat. Pass through the newer
    # `strategy` and `enabled` fields (both optional; render_usp.load_usps()
    # defaults strategy="" and enabled=True when absent, for backward compat
    # with pre-Revision-2 3-field USP JSON).
    normalized = []
    for u in usps_payload:
        normalized.append({
            "title": u.get("title") or u.get("headline"),
            "description": u.get("description") or u.get("support", ""),
            "proof": u.get("proof") or u.get("support", ""),
            "strategy": u.get("strategy", ""),
            "enabled": u.get("enabled", True),
        })
    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        theme=theme,
        usps=json.dumps(normalized),
        usps_json=None,
        wedge=inputs.get("wedge"),
        wedge_support=inputs.get("wedge_support"),
        out_dir=str(out_dir),
        language=language,
    )

    usps = _usp.load_usps(args)
    pptx_path = out_dir / f"{slug}_usp_{theme}.pptx"
    if theme == "dark":
        _usp.render_dark(args, usps, str(pptx_path))
    else:
        _usp.render_light(args, usps, str(pptx_path))
    return pptx_path


def render_reach(inputs: dict[str, Any], theme: str, out_dir: Path,
                  language: str = "en") -> Path:
    """Render Step 3 — How We Reach. Returns PPTX path."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    inner_choice = inputs.get("inner", "other")
    game_type = inputs.get("game_type", "custom")

    # reach payload from app form: list of 4 {cohort, channel, message, kpi}
    # skill expects: list of 4 {channels:[...], message, kpi}
    reach_payload = []
    for r in inputs["reach"]:
        ch = r.get("channels") or r.get("channel")
        if isinstance(ch, str):
            ch = [c.strip() for c in ch.split(",") if c.strip()]
        reach_payload.append({
            "channels": ch,
            "message": r["message"],
            "kpi": r["kpi"],
        })

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        type=game_type,
        inner=inner_choice,
        inner_name=inputs["cohorts"][0]["name"] if inner_choice == "other" else None,
        ring2_name=inputs["cohorts"][1]["name"] if game_type == "custom" else None,
        theme=theme,
        reach=json.dumps(reach_payload),
        reach_json=None,
        out_dir=str(out_dir),
        language=language,
    )

    reach = _reach.load_reach(args.reach, is_path=False)
    pptx_path = out_dir / f"{slug}_reach_{theme}.pptx"
    if theme == "dark":
        _reach.render_dark(args, reach, str(pptx_path))
    else:
        _reach.render_light(args, reach, str(pptx_path))
    return pptx_path


def render_roadmap(inputs: dict[str, Any], theme: str, release_date: dt.date,
                   out_dir: Path, phases_override: dict | None = None,
                   language: str = "en") -> Path:
    """Render Step 4 — Roadmap (6 sub-slides). Returns PPTX path."""
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    if phases_override is not None:
        phases = phases_override
    else:
        with open(ASSETS_DIR / "roadmap_phases.json") as f:
            phases = json.load(f)

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        release_date=release_date,
        theme=theme,
        phases_json=None,
        out_dir=str(out_dir),
        no_png=True,  # we'll convert via merged-deck PDF later
        language=language,
    )

    anchors = _roadmap.compute_anchors(release_date)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = phases["slides"]
    summary_events = phases.get("summary_events", {}).get("events", [])
    total = len(slides_data) + (1 if summary_events else 0)

    for i, sd in enumerate(slides_data):
        if theme == "dark":
            _roadmap.render_one_dark(prs, sd, i + 1, total, anchors, args)
        else:
            _roadmap.render_one_light(prs, sd, i + 1, total, anchors, args)

    if summary_events:
        slide_num = len(slides_data) + 1
        if theme == "dark":
            _roadmap.render_summary_dark(prs, summary_events, slide_num, total, anchors, args)
        else:
            _roadmap.render_summary_light(prs, summary_events, slide_num, total, anchors, args)

    pptx_path = out_dir / f"{slug}_roadmap_{theme}.pptx"
    prs.save(str(pptx_path))
    return pptx_path


def render_commercial_potential(inputs: dict[str, Any], theme: str, out_dir: Path,
                                  language: str = "en") -> Path:
    """Render Median Commercial Potential slide. Returns PPTX path.

    Output position 2 in the 12-slide pack (internal skill Step 5).

    Currency units (see gtm_revisions_summary.md for the correction history):
      - inputs["median_revenue_usd_millions"]: float, MILLIONS of dollars.
      - inputs["avg_price_usd"]: float, PLAIN dollars.
      - inputs["median_units_sold"]: int, raw unit count.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    platforms = inputs.get("platforms", ["PC", "PS5", "XSX", "SWITCH2"])
    if isinstance(platforms, str):
        platforms = [p.strip() for p in platforms.split(",") if p.strip()]

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        theme=theme,
        comp_set_name=inputs["comp_set_name"],
        median_revenue_usd_millions=float(inputs["median_revenue_usd_millions"]),
        median_units_sold=int(inputs["median_units_sold"]),
        avg_price_usd=float(inputs["avg_price_usd"]),
        avg_hours_played=float(inputs["avg_hours_played"]),
        platforms=",".join(platforms),
        out_dir=str(out_dir),
        language=language,
    )
    # render_commercial_potential.parse_args() derives comp_set_titles from
    # comp_set_name via regex (e.g. "Horror -- 19 titles" -> "19"). We call
    # render_dark/render_light directly (bypassing parse_args), so replicate
    # that derivation here -- required, not optional (see smoke-test note in
    # gtm_revisions_summary.md).
    m = re.search(r"(\d+)", args.comp_set_name)
    args.comp_set_titles = m.group(1) if m else "N"

    pptx_path = out_dir / f"{slug}_commercial_potential_{theme}.pptx"
    if theme == "dark":
        _commercial_potential.render_dark(args, str(pptx_path))
    else:
        _commercial_potential.render_light(args, str(pptx_path))
    return pptx_path


def render_commercial_risks(inputs: dict[str, Any], theme: str, out_dir: Path,
                             language: str = "en") -> Path:
    """Render Commercial Risks slide. Returns PPTX path.

    Output position 11 in the 12-slide pack (internal skill Step 6).
    inputs["risks"]: list of 1-5 {threat_level, proof, mitigation} dicts.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    # Prefer a Risks-specific wedge (inputs["risks_wedge"]) if the caller set
    # one, so the USP slide and Commercial Risks slide can carry different
    # wedge copy even though FormInputs shares one dict. Falls back to the
    # shared "wedge"/"wedge_support" keys (or the renderer's own default) if
    # no risks-specific override is present.
    risks_wedge = inputs.get("risks_wedge") or inputs.get("wedge")
    risks_wedge_support = inputs.get("risks_wedge_support") or inputs.get("wedge_support")

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        theme=theme,
        risks=json.dumps(inputs["risks"]),
        risks_json=None,
        wedge=risks_wedge,
        wedge_support=risks_wedge_support,
        out_dir=str(out_dir),
        language=language,
    )

    risks = _commercial_risks.load_risks(args)
    pptx_path = out_dir / f"{slug}_commercial_risks_{theme}.pptx"
    if theme == "dark":
        _commercial_risks.render_dark(args, risks, str(pptx_path))
    else:
        _commercial_risks.render_light(args, risks, str(pptx_path))
    return pptx_path


def render_description_razors(inputs: dict[str, Any], theme: str, out_dir: Path,
                               language: str = "en") -> Path:
    """Render Game Description & Razors slide. Returns PPTX path.

    Output position 12 (last slide) in the 12-slide pack (internal skill Step 7).
    inputs: accepts EITHER "description_100" (canonical key -- matches
    FormInputs/sample_inputs.py) OR "description" (renderer-script CLI arg
    name / legacy callers) for the ~100-word description. "description_100"
    is checked first. razor_20 (~20 words), razor_10 (~10 words).
    Word-count limits are soft warnings, not hard failures (except empty fields).
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    description = inputs.get("description_100")
    if description is None:
        description = inputs["description"]  # legacy/renderer-arg key fallback

    args = Namespace(
        title=inputs["title"],
        genre=inputs["genre"],
        theme=theme,
        description=description,
        razor_20=inputs["razor_20"],
        razor_10=inputs["razor_10"],
        out_dir=str(out_dir),
        language=language,
    )

    counts = _description_razors.validate_inputs(args)
    pptx_path = out_dir / f"{slug}_description_razors_{theme}.pptx"
    if theme == "dark":
        _description_razors.render_dark(args, counts, str(pptx_path))
    else:
        _description_razors.render_light(args, counts, str(pptx_path))
    return pptx_path


# ── Full-pack merger ─────────────────────────────────────────────────────────


def _copy_slide_from(target_prs: Presentation, source_prs: Presentation, slide_index: int):
    """Copy a slide from source_prs to target_prs at the end.

    python-pptx doesn't have a public slide-copy API. We use a low-level XML
    surgery approach: clone the slide XML and re-register the relationships.
    """
    from copy import deepcopy
    from pptx.oxml.ns import qn

    source_slide = source_prs.slides[slide_index]

    # Pick a blank layout on the target deck
    blank_layout = target_prs.slide_layouts[6]  # "Blank"
    new_slide = target_prs.slides.add_slide(blank_layout)

    # Remove placeholders that came with the blank layout
    for shp in list(new_slide.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    # Copy all shape XML from source slide
    for shp in source_slide.shapes:
        new_el = deepcopy(shp._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Copy slide-level background fill if present
    src_cSld = source_slide._element.find(qn("p:cSld"))
    tgt_cSld = new_slide._element.find(qn("p:cSld"))
    src_bg = src_cSld.find(qn("p:bg"))
    if src_bg is not None:
        # Replace target bg if it exists, otherwise insert at front
        tgt_bg = tgt_cSld.find(qn("p:bg"))
        if tgt_bg is not None:
            tgt_cSld.remove(tgt_bg)
        tgt_cSld.insert(0, deepcopy(src_bg))


def render_full_pack(
    inputs: dict[str, Any],
    theme: str,
    out_dir: Path,
    *,
    release_date: dt.date | None = None,
    phases_override: dict | None = None,
    language: str = "en",
) -> Path:
    """Render the complete 12-slide GTM pack.

    Returns the path to the merged PPTX containing all 12 slides, in OUTPUT
    order (internal skill step numbers differ -- see module docstring):
      1:     Sizing Circle
      2:     Median Commercial Potential   (NEW)
      3:     USP / Pillars
      4:     How We Reach
      5-10:  Roadmap (4.1 through 4.6)
      11:    Commercial Risks               (NEW)
      12:    Description & Razors           (NEW)

    inputs must include the fields required by every sub-renderer, notably:
      - median_revenue_usd_millions (float, MILLIONS of dollars)
      - avg_price_usd (float, plain dollars)
      - median_units_sold (int, raw count)
      - avg_hours_played (float)
      - comp_set_name (str)
      - platforms (list[str] or comma-separated str)
      - risks (list of 1-5 {threat_level, proof, mitigation})
      - description, razor_20, razor_10 (str)

    The release_date kwarg overrides inputs["release_date"] if provided.
    """
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    slug = _slugify(inputs["title"])

    # Resolve release date
    if release_date is None:
        rd = inputs.get("release_date")
        if rd is None:
            raise ValueError("release_date is required (in inputs or as kwarg)")
        release_date = dt.date.fromisoformat(rd) if isinstance(rd, str) else rd

    # Render the seven sub-decks in parallel (CPU-light, IO-heavy)
    tmp = Path(tempfile.mkdtemp(prefix="gtm_pack_"))
    try:
        with ThreadPoolExecutor(max_workers=7) as ex:
            futures = {
                "sizing":              ex.submit(render_sizing_circle, inputs, theme, tmp,
                                                 language),
                "commercial_potential": ex.submit(render_commercial_potential, inputs, theme,
                                                   tmp, language),
                "usp":                 ex.submit(render_usp, inputs, theme, tmp, language),
                "reach":               ex.submit(render_reach, inputs, theme, tmp, language),
                "roadmap":             ex.submit(render_roadmap, inputs, theme, release_date,
                                                 tmp, phases_override, language),
                "commercial_risks":    ex.submit(render_commercial_risks, inputs, theme, tmp,
                                                  language),
                "description_razors":  ex.submit(render_description_razors, inputs, theme,
                                                   tmp, language),
            }
            results = {k: f.result() for k, f in futures.items()}

        # Merge in FINAL OUTPUT order (locked, per user direction 2026-07-15):
        #   1. sizing
        #   2. commercial_potential
        #   3. usp
        #   4. commercial_risks       (moved up from position 11 on 2026-07-15)
        #   5. description_razors     (moved up from position 12 on 2026-07-15)
        #   6. reach
        #   7-12. roadmap (6 sub-slides)
        # Rationale for the reorder: put the strategic-positioning slides
        # (USPs → Risks → Razors) as a contiguous block so the deck reads
        # "who we are + what could go wrong + how we say it" before the
        # tactical execution slides (reach + roadmap).
        merged = Presentation()
        merged.slide_width = Inches(13.333)
        merged.slide_height = Inches(7.5)
        # Remove default slide that gets added with blank layout
        for order_key in ("sizing", "commercial_potential", "usp",
                          "commercial_risks", "description_razors",
                          "reach", "roadmap"):
            src = Presentation(str(results[order_key]))
            for i in range(len(src.slides)):
                _copy_slide_from(merged, src, i)

        merged_path = out_dir / f"{slug}_gtm_pack_{theme}.pptx"
        merged.save(str(merged_path))
        return merged_path
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ── Convenience: render + PDF + PNGs in one call ─────────────────────────────


def render_pack_with_artifacts(
    inputs: dict[str, Any],
    theme: str,
    out_dir: Path,
    *,
    release_date: dt.date | None = None,
    phases_override: dict | None = None,
    preview_dpi: int = 110,
    pdf_dpi: int = 150,
    skip_pngs: bool = False,
    language: str = "en",
) -> dict[str, Any]:
    """Render the full pack and produce PPTX, PDF, and per-slide PNGs.

    Returns:
        {
            "pptx": Path,
            "pdf": Path,
            "pngs": [Path, ...]   # one per slide, in order
        }
    """
    out_dir = Path(out_dir)
    pptx_path = render_full_pack(
        inputs, theme, out_dir,
        release_date=release_date,
        phases_override=phases_override,
        language=language,
    )
    pdf_path = _convert_to_pdf(pptx_path, out_dir)
    pngs = [] if skip_pngs else _pdf_to_pngs(pdf_path, out_dir, dpi=preview_dpi,
                                              prefix=pptx_path.stem)
    return {"pptx": pptx_path, "pdf": pdf_path, "pngs": pngs}


__all__ = [
    "render_sizing_circle",
    "render_commercial_potential",
    "render_usp",
    "render_reach",
    "render_roadmap",
    "render_commercial_risks",
    "render_description_razors",
    "render_full_pack",
    "render_pack_with_artifacts",
    "ASSETS_DIR",
]
