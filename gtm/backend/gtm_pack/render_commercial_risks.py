#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 6 'GTM Challenges' slide (Manifesto-derived
full-width risk list).

v8 polish pass (2026-07-18): text truncation is no longer acceptable on this
slide -- every mitigation/description word must be readable (these are
strategic decisions). The page-count rule is now DENSITY-AWARE instead of
purely count-based: we first try to fit all risks for a page at progressively
smaller font sizes (down to a 7pt floor for proof/mitigation, 10pt for the
threat pill/title-ish elements), and only if the floor rung still doesn't fit
do we split that page into more pages (re-chunking, never truncating). A
page's eyebrow gets a "(N OF M)" suffix whenever the deck ends up with more
than one page, exactly like the USP slide. Footer stays removed (global v7
rule).

Fallback order (locked): shrink fonts first -> then add pages -> only
truncate as an absolute last resort (single risk still doesn't fit at the
floor on its own page), and a RuntimeWarning is raised if that ever happens.

Threat levels: critical | high | medium | low (case-insensitive input,
rendered UPPERCASE). Color coding:
  Critical = red    (#D63A57 light / #E5615A dark)
  High     = orange (#E5A700 light / #FFB454 dark)
  Medium   = gold   (#FFC94D both themes)
  Low      = teal   (#1F9B8E light / #2FA9BD dark)

Outputs both a PPTX and a PNG to --out-dir. `render_dark()`/`render_light()`
add ONE OR MORE slides to the given Presentation and return it.
"""
from __future__ import annotations

import argparse
from ._title_fit import fit_title_pt
import json
import math
import os
import re
import subprocess
import tempfile

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def est_wrapped_lines(text: str, width_in: float, pt_size: float,
                       chars_per_inch_at_10pt: float = 16.5,
                       safety_margin: float = 1.18) -> int:
    """Estimate wrapped-line count. Same calibration as render_usp.py /
    render_reach.py -- deliberately over-counts so rows get slightly more
    height rather than risk overlap."""
    if not text:
        return 1
    chars_per_inch = chars_per_inch_at_10pt * (10.0 / pt_size) / safety_margin
    chars_per_line = max(8, int(width_in * chars_per_inch))
    total = 0
    for segment in text.split("\n"):
        n_chars = len(segment)
        total += max(1, math.ceil(n_chars / chars_per_line))
    return max(1, total)


def chunk_items(items: list, chunk_size: int = 3) -> list[list]:
    """Legacy count-based split rule, kept ONLY as the starting point for
    the density-aware re-chunker below (`paginate_risks`). Do not call this
    directly for final page counts any more -- see module docstring."""
    if len(items) <= chunk_size:
        return [items]
    return [items[:chunk_size], items[chunk_size:]]


def _page_fits(risks_page, area_h: float, *, text_w: float, min_gap: float = 0.10) -> bool:
    """True if `risks_page` fits within area_h at the floor font rung
    WITHOUT any truncation."""
    step = FONT_STEPS[-1]
    n = len(risks_page)
    heights = _measure_all(risks_page, step, text_w=text_w)
    total = sum(heights) + min_gap * (n - 1)
    return total <= area_h


def paginate_risks(risks: list, area_h: float, *, text_w: float,
                    max_per_page: int = 3) -> list[list]:
    """Density-aware pagination (v8 fix): start from the count-based split
    (<=3 per page) and, for any resulting page that would NOT fit even at
    the floor font rung, split that page further (greedily peeling items
    off the front) until every page fits without truncation. This is the
    "page more instead of truncating" behavior the design brief calls for:
    3 risks with dense mitigation copy on one slide is unrealistic, so we
    honestly split to 2 pages (or more) rather than clipping text.
    """
    coarse_pages = chunk_items(risks, chunk_size=max_per_page)
    final_pages: list[list] = []
    for page in coarse_pages:
        remaining = list(page)
        while remaining:
            # Try the largest prefix of `remaining` that fits at the floor
            # font rung; shrink the prefix size until it fits (min size 1).
            fit_size = len(remaining)
            while fit_size > 1 and not _page_fits(remaining[:fit_size], area_h, text_w=text_w):
                fit_size -= 1
            final_pages.append(remaining[:fit_size])
            remaining = remaining[fit_size:]
    return final_pages


# ---------- shape primitives ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w_pt=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w_pt)
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return box


# ---------- risk data shape ----------
VALID_LEVELS = ("critical", "high", "medium", "low")


def load_risks(args) -> list[tuple[str, str, str]]:
    """Return list of (threat_level_upper, proof, mitigation). Length 1-5."""
    if args.risks_json:
        with open(args.risks_json, "r") as f:
            raw = json.load(f)
    else:
        raw = json.loads(args.risks)
    if not isinstance(raw, list):
        raise SystemExit("--risks / --risks-json must be a JSON list")
    if not (1 <= len(raw) <= 5):
        raise SystemExit(f"Risk count must be 1-5; got {len(raw)}")
    out = []
    for i, item in enumerate(raw):
        if not isinstance(item, dict):
            raise SystemExit(f"Risk #{i+1} must be a JSON object")
        try:
            level = item["threat_level"].strip().lower()
            proof = item["proof"].strip()
            mitigation = item["mitigation"].strip()
        except (KeyError, AttributeError, TypeError):
            raise SystemExit(f"Risk #{i+1} missing one of: threat_level, proof, mitigation")
        if level not in VALID_LEVELS:
            raise SystemExit(
                f"Risk #{i+1} threat_level must be one of {VALID_LEVELS}; got '{level}'"
            )
        out.append((level.upper(), proof, mitigation))
    return out


def level_color(level_upper: str, theme: str) -> RGBColor:
    level = level_upper.lower()
    if theme == "dark":
        table = {
            "critical": hex_rgb("#E5615A"),
            "high":     hex_rgb("#FFB454"),
            "medium":   hex_rgb("#FFC94D"),
            "low":      hex_rgb("#2FA9BD"),
        }
    else:
        table = {
            "critical": hex_rgb("#D63A57"),
            "high":     hex_rgb("#E5A700"),
            "medium":   hex_rgb("#FFC94D"),
            "low":      hex_rgb("#1F9B8E"),
        }
    return table[level]


# ---------- Row height measurement (font-shrink ladder, mirrors render_usp.py) ----------
# v8 fix: extended down to a genuine 7pt floor for proof/mitigation (title
# pill text stays fixed at 9.5pt regardless of rung -- it's a short fixed
# label, not flowing body copy). Reaching this floor happens ONLY on a
# single-risk page with extremely dense copy; the density-aware pagination
# in `paginate_risks` is what normally keeps pages off the floor rungs.
FONT_STEPS = [
    dict(proof=13,   mitigation=11.5, proof_lh=0.24,  mit_lh=0.205),
    dict(proof=12,   mitigation=10.7, proof_lh=0.225, mit_lh=0.19),
    dict(proof=11,   mitigation=10.0, proof_lh=0.205, mit_lh=0.175),
    dict(proof=10.2, mitigation=9.3,  proof_lh=0.19,  mit_lh=0.16),
    dict(proof=9.0,  mitigation=8.2,  proof_lh=0.165, mit_lh=0.145),
    dict(proof=8.0,  mitigation=7.5,  proof_lh=0.145, mit_lh=0.135),
    dict(proof=7.5,  mitigation=7.0,  proof_lh=0.135, mit_lh=0.125),
]

ROW_PAD_TOP = 0.16
ROW_PAD_BOTTOM = 0.16
GAP_AFTER_PROOF = 0.10


def _measure_row(proof: str, mitigation: str, step, *, text_w: float) -> float:
    proof_lines = est_wrapped_lines(f"\u2192 {proof}", text_w, step["proof"])
    mit_lines = est_wrapped_lines(f"\u00bb {mitigation}", text_w, step["mitigation"])
    return (ROW_PAD_TOP
            + proof_lines * step["proof_lh"]
            + GAP_AFTER_PROOF
            + mit_lines * step["mit_lh"]
            + ROW_PAD_BOTTOM)


def _measure_all(risks_page, step, *, text_w) -> list[float]:
    return [_measure_row(proof, mitigation, step, text_w=text_w)
            for (_level, proof, mitigation) in risks_page]


def _fit_rows(risks_page, area_h: float, *, text_w: float, min_gap: float = 0.10):
    """Shrink-ladder fit -- NEVER truncates text under normal operation.
    Because `paginate_risks` already guarantees every page fits at the
    floor font rung before this function is ever called, the loop below
    almost always returns on an early (larger-font) rung. The only way to
    reach the post-loop fallback is a single-risk page whose copy is so
    dense it doesn't fit even at the 7pt floor -- an absolute edge case.
    In that case we scale row heights down to exactly fill the area (rows
    get tight but text is still fully rendered, just tightly spaced) and
    raise a RuntimeWarning so it's visible in logs; we do NOT clip text.
    Returns (risks_page, step, heights, gap)."""
    n = len(risks_page)
    for step in FONT_STEPS:
        heights = _measure_all(risks_page, step, text_w=text_w)
        total = sum(heights) + min_gap * (n - 1)
        if total <= area_h:
            leftover = area_h - total
            extra_gap = min(leftover / max(1, n - 1), 0.35) if n > 1 else 0.0
            remaining = leftover - extra_gap * (n - 1)
            total_h = sum(heights)
            heights = [h + (remaining * (h / total_h) if total_h > 0 else 0) for h in heights]
            gap = min_gap + extra_gap
            return risks_page, step, heights, gap

    # Floor rung still doesn't fit. Per the locked fallback order (shrink
    # -> page -> truncate LAST), this should be unreachable in practice
    # because paginate_risks already split pages down to a size that fits
    # at the floor. If it's ever hit anyway, do NOT clip text -- instead
    # compress the vertical rhythm (row gap + padding) as far as possible
    # while keeping every line of text rendered, and warn loudly.
    import warnings
    step = FONT_STEPS[-1]
    heights = _measure_all(risks_page, step, text_w=text_w)
    total = sum(heights) + min_gap * (n - 1)
    warnings.warn(
        f"render_commercial_risks: page with {n} risk(s) did not fit even "
        f"at the 7pt floor rung (needs {total:.2f}in, have {area_h:.2f}in). "
        "paginate_risks should have split this page further -- rendering "
        "with compressed spacing instead of truncating any text.",
        RuntimeWarning,
        stacklevel=2,
    )
    scale = area_h / total if total > 0 else 1.0
    heights = [h * scale for h in heights]
    gap = min_gap * scale
    return risks_page, step, heights, gap


def _render_page(slide, args, risks_page, page_num, total_pages, *, theme, colors, l_lang):
    BG, INK, MUTED, ACCENT = colors["BG"], colors["INK"], colors["MUTED"], colors["ACCENT"]
    LINE = colors["LINE"]

    if theme == "dark":
        add_rect(slide, 0, 0, 13.333, 7.5, BG)
        add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
        eyebrow_color = ACCENT
        title_x, title_y_e, title_y_t, title_y_s = 0.6, 0.4, 0.75, 1.55
        rx, ry, rw = 0.6, 2.35, 12.13
    else:
        add_rect(slide, 0, 0, 13.333, 7.5, BG)
        add_rect(slide, 0, 0, 0.25, 7.5, colors["C3"])
        eyebrow_color = colors["C3"]
        title_x, title_y_e, title_y_t, title_y_s = 0.7, 0.5, 0.85, 1.65
        rx, ry, rw = 0.7, 2.45, 12.0

    suffix = f" ({page_num} OF {total_pages})" if total_pages > 1 else ""
    add_text(slide, title_x, title_y_e, 10, 0.3, f"GTM CHALLENGES{suffix}",
             font="Trebuchet MS" if theme == "dark" else "Calibri",
             size=10 if theme == "dark" else body_pt(l_lang, 10),
             bold=True, color=eyebrow_color)
    add_text(slide, title_x, title_y_t, 12, 0.85, f"GTM Challenges \u2014 {args.title}",
             font="Trebuchet MS", size=fit_title_pt(f"GTM Challenges \u2014 {args.title}", 12), bold=True, color=INK)
    add_text(slide, title_x, title_y_s, 12, 0.4,
             f"{args.genre} \u00b7 Each challenge below is tracked with a named owner and a concrete mitigation.",
             font="Calibri", size=body_pt(l_lang, 13 if theme == "dark" else 14), color=MUTED)

    n = len(risks_page)
    area_h = 7.05 - ry  # no footer -- generous bottom margin
    pill_col_w = 1.05
    text_x = rx + pill_col_w
    text_w = rw - pill_col_w

    risks_page, step, heights, gap = _fit_rows(risks_page, area_h, text_w=text_w)

    y = ry
    for i, (level, proof, mitigation) in enumerate(risks_page):
        h = heights[i]
        c = level_color(level, theme)
        pill_w, pill_h = 0.9, 0.30
        pill = add_rect(slide, rx, y + ROW_PAD_TOP - 0.02, pill_w, pill_h, c)
        tf = pill.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = level
        r.font.name = "Trebuchet MS"
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = BG

        proof_lines = est_wrapped_lines(f"\u2192 {proof}", text_w, step["proof"])
        proof_h = proof_lines * step["proof_lh"] + 0.06
        add_text(slide, text_x, y + ROW_PAD_TOP, text_w, proof_h, f"\u2192 {proof}",
                 font="Calibri", size=body_pt(l_lang, step["proof"]), bold=True, color=c,
                 line_spacing=1.02)
        y_mit = y + ROW_PAD_TOP + proof_lines * step["proof_lh"] + GAP_AFTER_PROOF
        mit_lines = est_wrapped_lines(f"\u00bb {mitigation}", text_w, step["mitigation"])
        mit_h = mit_lines * step["mit_lh"] + 0.08
        add_text(slide, text_x, y_mit, text_w, mit_h, f"\u00bb {mitigation}",
                 font="Calibri", size=body_pt(l_lang, step["mitigation"]), color=INK,
                 line_spacing=1.05)

        if i < n - 1:
            add_rect(slide, rx, y + h - gap * 0.45, rw, 0.008, LINE)
        y += h + gap


# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, risks, out_path_or_prs):
    L = getattr(args, "language", "en")
    colors = dict(
        BG=hex_rgb("#0E1116"), LINE=hex_rgb("#1F2530"),
        INK=hex_rgb("#E8E6E1"), MUTED=hex_rgb("#8A8F99"),
        ACCENT=hex_rgb("#FFB454"),
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ry, rw, pill_col_w = 2.35, 12.13, 1.05
    area_h = 7.05 - ry
    text_w = rw - pill_col_w
    pages = paginate_risks(risks, area_h, text_w=text_w)
    total_pages = len(pages)
    for page_num, risks_page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_page(slide, args, risks_page, page_num, total_pages,
                     theme="dark", colors=colors, l_lang=L)

    prs.save(out_path_or_prs)
    return prs


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, risks, out_path_or_prs):
    L = getattr(args, "language", "en")
    colors = dict(
        BG=hex_rgb("#FFFFFF"), LINE=hex_rgb("#E8E8E8"),
        INK=hex_rgb("#1A1A1A"), MUTED=hex_rgb("#5C5C5C"),
        ACCENT=hex_rgb("#1F9B8E"), C3=hex_rgb("#1F9B8E"),
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ry, rw, pill_col_w = 2.45, 12.0, 1.05
    area_h = 7.05 - ry
    text_w = rw - pill_col_w
    pages = paginate_risks(risks, area_h, text_w=text_w)
    total_pages = len(pages)
    for page_num, risks_page in enumerate(pages, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_page(slide, args, risks_page, page_num, total_pages,
                     theme="light", colors=colors, l_lang=L)

    prs.save(out_path_or_prs)
    return prs


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_dir: str) -> list[str]:
    """Convert every slide/page in the PPTX to a PNG. Returns a list of PNG
    paths in slide order (one per page for multi-slide risk decks)."""
    os.makedirs(png_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = sorted(
            fn for fn in os.listdir(tmp) if fn.startswith("slide-") and fn.endswith(".png")
        )
        if not produced:
            raise RuntimeError("pdftoppm did not produce a PNG")
        out_paths = []
        for idx, fn in enumerate(produced, start=1):
            dst_path = os.path.join(png_dir, f"{base}_p{idx}.png")
            with open(os.path.join(tmp, fn), "rb") as src, open(dst_path, "wb") as dst:
                dst.write(src.read())
            out_paths.append(dst_path)
    return out_paths


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True)
    p.add_argument("--theme", required=True, choices=["dark", "light"])
    grp = p.add_mutually_exclusive_group(required=True)
    grp.add_argument("--risks", help="Inline JSON list of {threat_level, proof, mitigation} objects (1-5)")
    grp.add_argument("--risks-json", help="Path to a JSON file containing the risk list")
    p.add_argument("--wedge", default=None,
                   help="Optional manifesto statement, pipe-separated lines (e.g. 'A|B|C'). Unused in v7 full-width layout, kept for CLI compatibility.")
    p.add_argument("--wedge-support", default=None,
                   help="Optional one-line supporting sentence under the manifesto. Unused in v7 full-width layout, kept for CLI compatibility.")
    p.add_argument("--out-dir", required=True)
    return p.parse_args()


def main():
    args = parse_args()
    risks = load_risks(args)
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_commercial_risks_{args.theme}.pptx")
    if args.theme == "dark":
        render_dark(args, risks, pptx_path)
    else:
        render_light(args, risks, pptx_path)
    png_paths = convert_to_png(pptx_path, out_dir)
    print(f"PPTX: {pptx_path}")
    for pp in png_paths:
        print(f"PNG:  {pp}")


if __name__ == "__main__":
    main()
