#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 2 'USP / Pillars' slide (V2 Manifesto layout).

Two themes, parity with Step 1:
  - dark   (V2 Modern Mono)  : dark slide, teal ramp + warm gold accent on last pillar
  - light  (V4 Bold Brand)   : light slide, mint/teal/rose/gold tier ramp, left teal stripe

v7 polish pass (2026-07-18): the single fixed-height full-width row layout
that caused severe text-on-text overlap for longer descriptions/proof lines
is replaced with a CARD layout (bordered card per USP, generous internal
padding) plus a MULTI-SLIDE split: 1-3 enabled USPs render on one slide;
4-5 enabled USPs render across TWO slides (items 1-3, then remaining items),
each page carrying a "(N OF M)" eyebrow suffix. Card heights are computed
from an estimated wrapped-line count for the description text so long
descriptions get taller cards instead of overlapping the next card.
Footer removed (global v7 rule).

v7.2 polish pass (2026-07-18): fixes a "mess with text spacing" report on
this specific slide (Slide 3 in the assembled pack):
  - Card internal padding is now a fixed set of tokens (PAD_LEFT/RIGHT/
    TOP/BOTTOM) instead of a single reused `pad` value, so the description
    text always starts >= 0.35in from the card's left edge (well clear of
    the 0.06in tier stripe) and always leaves >= 0.4in of right padding
    before the card's right border.
  - Card heights on a page are now made UNIFORM: every card on the page
    takes the max natural content height (bounded by the available area),
    instead of each card being exactly as tall as its own content. A
    lopsided page (one huge card next to two short ones) reads as broken;
    equal-height cards with breathing room read as intentional.
  - Per-card titles now reuse `fit_title_pt` (the same helper used for the
    slide headline) so a long Card 02-style title auto-shrinks to fit on
    one line inside the card's text column instead of risking a collision
    with the card's right edge.
  - Vertical rhythm inside a card is now driven by fixed spacing tokens:
    always GAP_TITLE_DESC between title bottom and description top, always
    GAP_DESC_PROOF between description bottom and the proof line's hairline
    divider, always PAD_BOTTOM from the last line of text to the card's
    bottom edge.
  - A subtle 1px hairline divider is drawn directly above the proof line
    (when a proof line is present) to visually separate it from the
    description instead of letting it run straight into the body copy.
  - Description line spacing raised from a tight single-line default to
    1.25 for readability.

Outputs both a PPTX and a PNG to --out-dir. `render_dark()`/`render_light()`
now add ONE OR TWO slides to the given Presentation and return it.
"""
from __future__ import annotations

import argparse
import json
import math
import os
import re
import subprocess
import tempfile

from pptx import Presentation
from ._title_fit import fit_title_pt
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def est_wrapped_lines(text: str, width_in: float, pt_size: float,
                       chars_per_inch_at_10pt: float = 16.5,
                       safety_margin: float = 1.18) -> int:
    """Estimate how many lines `text` will wrap to inside a box of
    `width_in` inches at font size `pt_size`.

    Calibrated against real soffice-rendered Calibri output (v7 polish pass):
    a 10.5pt paragraph in a ~10.6in-wide box wraps at ~155 chars/line, i.e.
    ~14.6 chars/inch. `chars_per_inch_at_10pt` is scaled by 10/pt_size to
    approximate other sizes, then a `safety_margin` shrinks the effective
    chars-per-line so this estimator deliberately OVER-counts lines rather
    than under-counts -- under-counting is what caused the v6 overlap bug
    where proof/strategy text collided with a description that wrapped to
    more lines than budgeted.
    """
    if not text:
        return 1
    chars_per_inch = chars_per_inch_at_10pt * (10.0 / pt_size) / safety_margin
    chars_per_line = max(8, int(width_in * chars_per_inch))
    total = 0
    for segment in text.split("\n"):
        n_chars = len(segment)
        total += max(1, math.ceil(n_chars / chars_per_line))
    return max(1, total)


# ---------- shape primitives ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w_pt=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w_pt)
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return box


# ---------- USP data shape ----------
def load_usps(args) -> list[tuple[str, str, str, str]]:
    """Return list of (title, description, proof, strategy). Length 1-5.

    Backward compatible: entries without `enabled` are treated as enabled;
    entries without `strategy` get an empty strategy string (no strategy
    line rendered). Entries with enabled=false are dropped entirely before
    the 1-5 count check.
    """
    if args.usps_json:
        with open(args.usps_json, "r") as f:
            raw = json.load(f)
    else:
        raw = json.loads(args.usps)
    if not isinstance(raw, list):
        raise SystemExit("--usps / --usps-json must be a JSON list")

    enabled_items = []
    for i, item in enumerate(raw):
        if not isinstance(item, dict):
            raise SystemExit(f"USP #{i+1} must be a JSON object")
        if not item.get("enabled", True):
            continue  # skip disabled rows entirely
        try:
            t = item["title"].strip()
            d = item["description"].strip()
            p = item["proof"].strip()
        except (KeyError, AttributeError, TypeError):
            raise SystemExit(f"USP #{i+1} missing one of: title, description, proof")
        s = (item.get("strategy") or "").strip()
        enabled_items.append((t, d, p, s))

    if not (1 <= len(enabled_items) <= 5):
        raise SystemExit(f"Enabled USP count must be 1-5; got {len(enabled_items)}")
    return enabled_items


# ---------- accent ramps (4 entries; we slice for 3 or 5) ----------
def dark_accents() -> list[RGBColor]:
    A2 = hex_rgb("#155966")
    A3 = hex_rgb("#2FA9BD")
    A4 = hex_rgb("#7FD8E3")
    GOLD = hex_rgb("#FFB454")
    return [A4, A3, A2, GOLD]


def light_accents() -> list[RGBColor]:
    C4 = hex_rgb("#7DD4C9")  # mint
    C3 = hex_rgb("#1F9B8E")  # teal
    C2 = hex_rgb("#D63A57")  # rose
    C1 = hex_rgb("#E5A700")  # gold
    return [C4, C3, C2, C1]


def expand_accents(base: list[RGBColor], n: int) -> list[RGBColor]:
    """Expand/contract the 4-color ramp to exactly n colors (1-5 supported)."""
    if n == 1:
        return [base[0]]
    if n == 2:
        return [base[0], base[3]]
    if n == 3:
        return [base[0], base[1], base[3]]
    if n == 4:
        return list(base)
    if n == 5:
        mid = RGBColor(
            (base[1][0] + base[2][0]) // 2,
            (base[1][1] + base[2][1]) // 2,
            (base[1][2] + base[2][2]) // 2,
        )
        return [base[0], base[1], mid, base[2], base[3]]
    raise SystemExit(f"Unsupported USP count: {n}")


# ---------- multi-slide chunking (LOCKED split rule -- count-based, AUTHORITATIVE) ----------
def chunk_items(items: list) -> list[list]:
    """Split into pages purely by COUNT, per the user's original directive
    (locked verbatim this session): "1-3 items = 1 slide, 4-5 items = 2
    slides". This is now the ONLY thing that decides page count -- there is
    no density-based fallback that can add pages. Extending the same rule
    upward for larger sets: 6-8 items = 3 pages.

      1-3 items -> 1 page                      [n]
      4-5 items -> 2 pages, first gets 3       4 -> [3, 1], 5 -> [3, 2]
      6-8 items -> 3 pages, first two get 3    6 -> [3, 3], 7 -> [3, 3, 1],
                                                8 -> [3, 3, 2]

    If content is dense, the shrink ladder in `_render_page` is responsible
    for fitting it onto the page count decided here (shrinking fonts more
    aggressively, then truncating text as a last resort) -- it must NEVER
    ask for more pages than this function returns.
    """
    n = len(items)
    if n <= 3:
        return [items]
    if n <= 5:
        return [items[:3], items[3:]]
    # 6-8: two pages of 3 plus whatever remains (1 or 2)
    pages = []
    i = 0
    while i < n:
        pages.append(items[i:i + 3])
        i += 3
    return pages


FONT_STEPS_FLOOR = dict(title=11.0, desc=8.5, band=8.5, title_lh=0.22, desc_lh=0.145, band_lh=0.185)


# ---------- shared per-theme card renderer ----------
# ---------- shared per-theme card renderer ----------
# ---------- spacing tokens (v7.2 -- LOCKED rhythm) ----------
PAD_LEFT = 0.35     # from card left edge to text column left edge (clear of 0.06in stripe)
PAD_RIGHT = 0.40    # from text column right edge to card right edge
PAD_TOP = 0.28      # from card top edge to title top
PAD_BOTTOM = 0.40   # from last line of text (proof/strategy) to card bottom edge
GAP_TITLE_DESC = 0.18   # title bottom -> description top
GAP_DESC_PROOF = 0.22   # description bottom -> proof hairline/line top
GAP_PROOF_STRATEGY = 0.08  # proof bottom -> strategy top
DESC_LINE_SPACING = 1.25   # readability bump from the old ~1.08 default
NUM_COL_W = 0.6
NUM_TO_TEXT_GAP = 0.28  # gap between the "0N" numeral and the text column


def _measure_page(usps_page, text_w, desc_size, title_size,
                   title_line_h, desc_line_h, band_line_h,
                   band_text_size=None, *, pad_top=None, pad_bottom=None,
                   gap_title_desc=None, gap_desc_proof=None,
                   gap_proof_strategy=None):
    """Return (natural_heights, desc_lines_list, title_lines_list,
    proof_lines_list, strategy_lines_list) for the given font sizes --
    pure measurement, no drawing. Uses the locked spacing tokens
    (PAD_TOP/PAD_BOTTOM/GAP_*) for vertical rhythm so the measured height
    always matches what _render_page actually draws, UNLESS one of the
    optional pad_*/gap_* overrides is passed -- the floor rungs of the
    font shrink ladder also tighten these spacing tokens (not just font
    size) so 3-per-page dense content has a real chance of fitting instead
    of hitting the fixed per-card padding/gap overhead as a hard wall.

    Proof and strategy lines are each measured with est_wrapped_lines too
    (not assumed to be a single line) -- a long `\u2192 proof` line that wraps
    to 2 lines used to be budgeted as 1 line of height, which let it
    overflow past the card's bottom border. The `\u2192 `/`\u00bb ` prefix eats a
    little width, so we measure against a very slightly narrower box.
    """
    band_size = band_text_size if band_text_size is not None else desc_size
    band_text_w = max(1.0, text_w - 0.18)
    p_top = pad_top if pad_top is not None else PAD_TOP
    p_bottom = pad_bottom if pad_bottom is not None else PAD_BOTTOM
    g_td = gap_title_desc if gap_title_desc is not None else GAP_TITLE_DESC
    g_dp = gap_desc_proof if gap_desc_proof is not None else GAP_DESC_PROOF
    g_ps = gap_proof_strategy if gap_proof_strategy is not None else GAP_PROOF_STRATEGY
    natural_heights, desc_lines_list, title_lines_list = [], [], []
    proof_lines_list, strategy_lines_list = [], []
    for (title, desc, proof, strategy) in usps_page:
        title_lines = est_wrapped_lines(title, text_w, title_size)
        desc_lines = est_wrapped_lines(desc, text_w, desc_size) if desc else 0
        proof_lines = est_wrapped_lines(proof, band_text_w, band_size) if proof else 0
        strategy_lines = est_wrapped_lines(strategy, band_text_w, band_size) if strategy else 0
        content_h = p_top + max(1, title_lines) * title_line_h
        if desc_lines:
            content_h += g_td + desc_lines * desc_line_h
        if proof_lines or strategy_lines:
            content_h += g_dp
            if proof_lines:
                content_h += proof_lines * band_line_h
            if strategy_lines:
                content_h += g_ps + strategy_lines * band_line_h
        content_h += p_bottom
        natural_heights.append(content_h)
        desc_lines_list.append(desc_lines)
        title_lines_list.append(title_lines)
        proof_lines_list.append(proof_lines)
        strategy_lines_list.append(strategy_lines)
    return natural_heights, desc_lines_list, title_lines_list, proof_lines_list, strategy_lines_list


def _render_page(slide, args, usps_page, accents_page, page_num, total_pages,
                  *, theme: str, colors: dict):
    """Render one page (slide) of up to 3 USP cards. `colors` carries the
    theme-specific palette so this single function drives both themes.

    Layout strategy (v7.2 polish pass):
      1. Measure natural content height per card at the BASE font sizes
         (title 15pt, desc 10.5pt, proof/strategy ~10-10.5pt) using fixed
         spacing tokens for vertical rhythm (title->desc, desc->proof,
         bottom padding) so every card shares the same internal cadence.
      2. If the tallest card's natural height times the card count fits
         within the available area, apply that SAME height to every card
         on the page (uniform rows) and center each card's content
         vertically inside the extra room -- this replaces the old
         per-card variable-height layout that made a page with one long
         card and two short ones look lopsided.
      3. If it does NOT fit, shrink title/desc/band font sizes together in
         steps (down to a floor) and re-measure at each step until content
         fits, or the floor is reached.
      4. Only if even the floor size still doesn't fit (extreme edge case:
         very long text on all 3 cards) do we fall back to proportionally
         scaling the uniform card height down to fill the area exactly --
         this can leave text slightly tight but keeps cards from colliding.
    """
    L = getattr(args, "language", "en")
    BG, SURFACE, BORDER, INK, MUTED, ACCENT, STRIPE = (
        colors["BG"], colors["SURFACE"], colors["BORDER"], colors["INK"],
        colors["MUTED"], colors["ACCENT"], colors["STRIPE"],
    )
    HAIR = colors.get("HAIR", BORDER)

    if theme == "dark":
        add_rect(slide, 0, 0, 13.333, 7.5, BG)
        add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)
        margin_x = 0.6
        eyebrow_color = ACCENT
    else:
        add_rect(slide, 0, 0, 13.333, 7.5, BG)
        add_rect(slide, 0, 0, 0.25, 7.5, STRIPE)
        margin_x = 0.7
        eyebrow_color = STRIPE

    suffix = f" ({page_num} OF {total_pages})" if total_pages > 1 else ""
    eyebrow = f"USPs{suffix}"
    add_text(slide, margin_x, 0.4 if theme == "dark" else 0.5,
             10, 0.3, eyebrow,
             font=("Trebuchet MS" if theme == "dark" else "Calibri"),
             size=10 if theme == "dark" else body_pt(L, 10),
             bold=True, color=eyebrow_color)
    add_text(slide, margin_x, 0.75 if theme == "dark" else 0.85,
             12, 0.85, f"What sets {args.title} apart",
             font="Trebuchet MS", size=fit_title_pt(f"What sets {args.title} apart", 12), bold=True, color=INK)
    add_text(slide, margin_x, 1.55 if theme == "dark" else 1.65,
             12, 0.4,
             f"{args.genre} \u00b7 The pillars carrying the launch \u00b7 Each pillar is independently defensible.",
             font="Calibri", size=body_pt(L, 13 if theme == "dark" else 14), color=MUTED)

    # ---- Card list geometry ----
    rx = margin_x
    ry = 2.3 if theme == "dark" else 2.4
    rw = 13.333 - 2 * margin_x
    n = len(usps_page)
    area_top = ry
    area_bottom = 7.25
    area_h = area_bottom - area_top
    gap = 0.2
    text_x_offset = NUM_TO_TEXT_GAP + NUM_COL_W + PAD_LEFT - 0.28
    # text column width: from text start to card right edge, minus PAD_RIGHT
    text_w = rw - text_x_offset - PAD_RIGHT

    # ---- Font-size shrink ladder (title, desc, proof/strategy, line-heights) ----
    # Each rung shrinks together so type hierarchy is preserved.
    FONT_STEPS = [
        dict(title=15, desc=10.5, band=10.5, title_lh=0.32, desc_lh=0.185, band_lh=0.23,
             pad_top=PAD_TOP, pad_bottom=PAD_BOTTOM, gap_td=GAP_TITLE_DESC,
             gap_dp=GAP_DESC_PROOF, gap_ps=GAP_PROOF_STRATEGY),
        dict(title=14, desc=9.7,  band=9.7,  title_lh=0.30, desc_lh=0.175, band_lh=0.22,
             pad_top=PAD_TOP, pad_bottom=PAD_BOTTOM, gap_td=GAP_TITLE_DESC,
             gap_dp=GAP_DESC_PROOF, gap_ps=GAP_PROOF_STRATEGY),
        dict(title=13, desc=9.0,  band=9.0,  title_lh=0.28, desc_lh=0.165, band_lh=0.21,
             pad_top=PAD_TOP, pad_bottom=PAD_BOTTOM, gap_td=GAP_TITLE_DESC,
             gap_dp=GAP_DESC_PROOF, gap_ps=GAP_PROOF_STRATEGY),
        dict(title=12.5, desc=8.5, band=8.5, title_lh=0.26, desc_lh=0.155, band_lh=0.20,
             pad_top=PAD_TOP, pad_bottom=PAD_BOTTOM, gap_td=GAP_TITLE_DESC,
             gap_dp=GAP_DESC_PROOF, gap_ps=GAP_PROOF_STRATEGY),
        # v7.3 polish pass: the count-based split rule (1-3 items -> ALWAYS
        # 1 page) is now authoritative -- there is no density-based page-
        # count fallback any more, so the shrink ladder must be able to
        # reach a genuinely smaller floor before we ever truncate text.
        # These two extra rungs are the "shrink MORE aggressively before
        # truncating" step called for by the fix: title drops as low as
        # 11pt, description/proof/strategy as low as 8.5pt (already the
        # old floor) with tighter line-heights to match. Critically, these
        # floor rungs ALSO tighten the fixed padding/gap tokens (not just
        # font size) -- on a dense 3-per-page layout the fixed PAD_TOP/
        # PAD_BOTTOM/GAP_* overhead alone (title->desc gap, desc->proof
        # gap, proof->strategy gap, top/bottom padding) can eat most of a
        # card's tiny per-card budget before a single line of text is even
        # drawn, so the padding itself has to shrink at the floor or the
        # very shortest possible render still won't fit.
        dict(title=12.0, desc=8.5, band=8.5, title_lh=0.24, desc_lh=0.148, band_lh=0.19,
             pad_top=0.20, pad_bottom=0.28, gap_td=0.13, gap_dp=0.16, gap_ps=0.05),
        dict(title=11.0, desc=8.5, band=8.5, title_lh=0.22, desc_lh=0.145, band_lh=0.185,
             pad_top=0.14, pad_bottom=0.18, gap_td=0.09, gap_dp=0.11, gap_ps=0.04),
    ]

    chosen = FONT_STEPS[-1]
    natural_heights = desc_lines_list = title_lines_list = None
    proof_lines_list = strategy_lines_list = None
    for step in FONT_STEPS:
        # Per-card title auto-shrink: reuse fit_title_pt so a long title
        # (e.g. Card 02's full-width headline) shrinks to fit ONE line
        # inside the card's own text column width instead of overflowing
        # toward the card's right edge. fit_title_pt's own floor (22pt) is
        # too large relative to our card title base sizes, so we clamp the
        # per-card result between a card-appropriate floor and this rung's
        # base title size.
        title_sizes_this_step = [
            min(step["title"], max(9.5, fit_title_pt(t, text_w, base_pt=int(step["title"]), min_pt=9)))
            for (t, _, _, _) in usps_page
        ]
        nh, dl, tl, pl, sl = [], [], [], [], []
        for (title, desc, proof, strategy), tsize in zip(usps_page, title_sizes_this_step):
            _nh, _dl, _tl, _pl, _sl = _measure_page(
                [(title, desc, proof, strategy)], text_w, step["desc"], tsize,
                step["title_lh"], step["desc_lh"], step["band_lh"], step["band"],
                pad_top=step["pad_top"], pad_bottom=step["pad_bottom"],
                gap_title_desc=step["gap_td"], gap_desc_proof=step["gap_dp"],
                gap_proof_strategy=step["gap_ps"],
            )
            nh.append(_nh[0]); dl.append(_dl[0]); tl.append(_tl[0])
            pl.append(_pl[0]); sl.append(_sl[0])
        total_natural = sum(nh) + gap * (n - 1)
        if total_natural <= area_h:
            chosen = step
            natural_heights, desc_lines_list, title_lines_list = nh, dl, tl
            proof_lines_list, strategy_lines_list = pl, sl
            title_sizes = title_sizes_this_step
            break
    else:
        # Even the smallest font rung doesn't fit as-is. Rather than only
        # shrinking box heights (which left font size untouched and caused
        # text to spill past its box in extreme worst-case density), TRUNCATE
        # each description AND each proof/strategy line to a max-lines
        # budget at the floor font size. This is a last resort -- the
        # count-based split rule (locked this session) is authoritative
        # about page count, so a dense page must never turn into more
        # pages; it must instead lose some text.
        #
        # v7.3: the per-card line-cap search is now INDEPENDENT per card
        # instead of a single shared cap applied uniformly to every card
        # on the page. The old shared-cap version tightened every card's
        # cap in lockstep until the WORST (densest) card fit, which meant
        # short/normal cards got truncated just as hard as the one dense
        # card even though they had plenty of room -- e.g. a 3-card page
        # where only card 2 has both a long description and a proof line
        # would chop cards 1 and 3 down to a single line each for no
        # reason. Now each card searches its own line caps independently
        # against its OWN uniform per-card budget, so short cards keep
        # their full text and only the genuinely dense card(s) lose lines.
        chosen = FONT_STEPS[-1]
        band_text_w = max(1.0, text_w - 0.18)
        per_card_budget = (area_h - gap * (n - 1)) / n

        import warnings
        warnings.warn(
            f"render_usp: page with {n} card(s) did not fit even at the "
            f"floor font rung (title={chosen['title']}pt / desc={chosen['desc']}pt). "
            "Per the count-based split rule this page count is NOT negotiable -- "
            "truncating proof/strategy/description text instead of adding a page.",
            RuntimeWarning,
            stacklevel=2,
        )

        def _truncate_at_caps(desc, proof, strategy, max_desc_lines, max_band_lines):
            desc_budget_chars = int(text_w * 16.5 * (10.0 / chosen["desc"]) / 1.18) * max_desc_lines
            d = desc
            if d and len(d) > desc_budget_chars:
                d = d[: max(0, desc_budget_chars - 1)].rstrip() + "\u2026"
            band_budget_chars = int(band_text_w * 16.5 * (10.0 / chosen["band"]) / 1.18) * max_band_lines
            pr = proof
            if pr and len(pr) > band_budget_chars:
                pr = pr[: max(0, band_budget_chars - 1)].rstrip() + "\u2026"
            st = strategy
            if st and len(st) > band_budget_chars:
                st = st[: max(0, band_budget_chars - 1)].rstrip() + "\u2026"
            return d, pr, st

        truncated_page, title_sizes_try = [], []
        nh, dl, tl, pl, sl = [], [], [], [], []
        for (title, desc, proof, strategy) in usps_page:
            tsize = min(chosen["title"], max(9.5, fit_title_pt(title, text_w, base_pt=int(chosen["title"]), min_pt=9)))
            max_desc_lines, max_band_lines = 6, 3  # generous starting caps, per-card
            d, pr, st = desc, proof, strategy
            card_nh = card_dl = card_tl = card_pl = card_sl = None
            while True:
                d, pr, st = _truncate_at_caps(desc, proof, strategy, max_desc_lines, max_band_lines)
                _nh, _dl, _tl, _pl, _sl = _measure_page(
                    [(title, d, pr, st)], text_w, chosen["desc"], tsize,
                    chosen["title_lh"], chosen["desc_lh"], chosen["band_lh"], chosen["band"],
                    pad_top=chosen["pad_top"], pad_bottom=chosen["pad_bottom"],
                    gap_title_desc=chosen["gap_td"], gap_desc_proof=chosen["gap_dp"],
                    gap_proof_strategy=chosen["gap_ps"],
                )
                card_nh, card_dl, card_tl = _nh[0], _dl[0], _tl[0]
                card_pl, card_sl = min(_pl[0], max_band_lines), min(_sl[0], max_band_lines)
                if card_nh <= per_card_budget or (max_desc_lines <= 1 and max_band_lines <= 1):
                    break
                # Tighten this card's own caps a notch and re-measure. Shrink
                # whichever budget (description vs. proof/strategy) is the
                # larger share of overflow first so short proof lines aren't
                # clipped before a still-generous description gets trimmed.
                if max_desc_lines > 2:
                    max_desc_lines -= 1
                elif max_band_lines > 1:
                    max_band_lines -= 1
                else:
                    max_desc_lines = max(1, max_desc_lines - 1)
            truncated_page.append((title, d, pr, st))
            title_sizes_try.append(tsize)
            nh.append(card_nh); dl.append(min(card_dl, max_desc_lines)); tl.append(card_tl)
            pl.append(card_pl); sl.append(card_sl)

        usps_page = truncated_page
        title_sizes = title_sizes_try
        natural_heights, desc_lines_list, title_lines_list = nh, dl, tl
        proof_lines_list, strategy_lines_list = pl, sl

    desc_size = chosen["desc"]
    band_size = chosen["band"]
    title_line_h = chosen["title_lh"]
    desc_line_h = chosen["desc_lh"]
    band_line_h = chosen["band_lh"]
    # Spacing tokens for THIS rung -- the floor rungs of the shrink ladder
    # tighten these too (not just font size), so the draw loop below must
    # use the SAME chosen pad/gap values that were used to measure card
    # heights, or drawn text could drift out of sync with its box again.
    pad_top = chosen.get("pad_top", PAD_TOP)
    gap_td = chosen.get("gap_td", GAP_TITLE_DESC)
    gap_dp = chosen.get("gap_dp", GAP_DESC_PROOF)
    gap_ps = chosen.get("gap_ps", GAP_PROOF_STRATEGY)

    # ---- UNIFORM card heights: every card on this page gets the SAME
    # height (the tallest card's natural content height), so a page never
    # looks lopsided with one giant card next to short ones. When the
    # for/else loop above already found a font rung that fits (the common
    # case), max_natural is comfortably <= per_card_budget and we simply
    # spread any leftover space as extra breathing room. In the rare
    # iterative-truncation fallback, max_natural is already <=
    # per_card_budget BY CONSTRUCTION (that loop only exits once it is),
    # so this branch never needs to scale text-bearing boxes below what
    # the text actually needs. ----
    max_natural = max(natural_heights)
    total_uniform = max_natural * n + gap * (n - 1)
    if total_uniform <= area_h:
        leftover = area_h - total_uniform
        if n == 1:
            # A lone card on a split-off page (the dense-content escape
            # hatch can leave a single remainder card by itself) would
            # otherwise sit at its bare natural height with a large empty
            # gap below it. Grow its box to use a healthy share of the
            # leftover room instead, capped so it doesn't look absurdly
            # stretched -- content stays top-aligned inside the taller box.
            uniform_h = max_natural + min(leftover * 0.6, 1.6)
            eff_gap = gap
        else:
            extra_gap = min(leftover / max(1, n), 0.5)
            uniform_h = max_natural
            eff_gap = gap + extra_gap
    else:
        # Should not happen given the guarantees above, but keep a safe
        # fallback: scale down uniformly so cards never overlap.
        scale = area_h / total_uniform if total_uniform > 0 else 1.0
        uniform_h = max_natural * scale
        eff_gap = gap * scale
    card_heights = [uniform_h] * n

    # ---- Render using computed heights + chosen font sizes ----
    y = area_top
    for i, ((title, desc, proof, strategy), c) in enumerate(zip(usps_page, accents_page)):
        card_h = card_heights[i]
        add_rect(slide, rx, y, rw, card_h, SURFACE if theme == "dark" else BG,
                 line=(BORDER if theme == "dark" else colors["HAIR"]), line_w_pt=1.0)
        add_rect(slide, rx, y, 0.06, card_h, c)

        text_x = rx + text_x_offset
        num_x = rx + PAD_LEFT

        title_size = title_sizes[i]
        add_text(slide, num_x, y + pad_top - 0.02, NUM_COL_W, 0.4, f"0{i+1}",
                 font="Trebuchet MS", size=16, bold=True, color=c)
        add_text(slide, text_x, y + pad_top - 0.02, text_w, 0.4, title,
                 font="Trebuchet MS", size=title_size, bold=True, color=INK)

        title_lines = title_lines_list[i]
        desc_y = y + pad_top + max(1, title_lines) * title_line_h + gap_td
        desc_lines = desc_lines_list[i]
        desc_h = max(0.2, desc_lines * desc_line_h + 0.08)
        if desc:
            add_text(slide, text_x, desc_y, text_w, desc_h, desc,
                     font="Calibri", size=body_pt(L, desc_size), color=MUTED,
                     line_spacing=DESC_LINE_SPACING)

        # Proof + strategy sit directly below the description's OWN measured
        # height (not a fixed footer band), so they never collide with
        # description text regardless of how many lines it wrapped to. Each
        # line's own box height now reflects its OWN wrapped-line count
        # (proof_lines_list[i] / strategy_lines_list[i]) instead of a fixed
        # single-line assumption, so a proof line that wraps to 2 lines
        # can never spill past the card's bottom border. A subtle hairline
        # divider is drawn just above the proof line to give it visual
        # separation from the description body copy.
        proof_lines = proof_lines_list[i]
        strategy_lines = strategy_lines_list[i]
        band_y = desc_y + desc_h + gap_dp
        if proof:
            hairline_y = band_y - gap_dp / 2
            add_rect(slide, text_x, hairline_y, text_w, 0.010, HAIR)
            proof_h = max(0.2, proof_lines * band_line_h + 0.04)
            add_text(slide, text_x, band_y, text_w, proof_h, f"\u2192 {proof}",
                     font="Calibri", size=body_pt(L, band_size), bold=True, color=c,
                     line_spacing=1.05)
            band_y += proof_h
        if strategy:
            band_y += gap_ps if proof else 0
            strategy_h = max(0.2, strategy_lines * band_line_h + 0.04)
            add_text(slide, text_x, band_y, text_w, strategy_h, f"\u00bb {strategy}",
                     font="Calibri", size=body_pt(L, band_size - 0.5), color=INK,
                     line_spacing=1.05)

        y += card_h + eff_gap



# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================

# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, usps, out_path_or_prs):
    """Render 1 or 2 slides (per the locked split rule) into a Presentation
    and save it. Accepts either a path (creates a fresh Presentation) so the
    CLI/standalone use case keeps working, matching prior call signature.
    """
    colors = dict(
        BG=hex_rgb("#0E1116"), SURFACE=hex_rgb("#161A21"), BORDER=hex_rgb("#1F2530"),
        INK=hex_rgb("#E8E6E1"), MUTED=hex_rgb("#8A8F99"), ACCENT=hex_rgb("#FFB454"),
        STRIPE=hex_rgb("#FFB454"), HAIR=hex_rgb("#1F2530"),
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    accents_all = expand_accents(dark_accents(), len(usps))
    # Page count is decided ONLY by the count-based split rule (locked
    # this session, verbatim): 1-3 items -> 1 page, 4-5 items -> 2 pages,
    # 6-8 items -> 3 pages. There is no density-based fallback that can
    # add pages -- if a page's content is dense, `_render_page`'s own
    # shrink ladder (and, as a last resort, text truncation) absorbs it.
    pages = chunk_items(list(zip(usps, accents_all)))
    total_pages = len(pages)

    for page_num, page in enumerate(pages, start=1):
        usps_page = [p[0] for p in page]
        accents_page = [p[1] for p in page]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_page(slide, args, usps_page, accents_page, page_num, total_pages,
                     theme="dark", colors=colors)

    prs.save(out_path_or_prs)
    return prs


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, usps, out_path_or_prs):
    colors = dict(
        BG=hex_rgb("#FFFFFF"), SURFACE=hex_rgb("#FFFFFF"), BORDER=hex_rgb("#E8E8E8"),
        INK=hex_rgb("#1A1A1A"), MUTED=hex_rgb("#5C5C5C"), ACCENT=hex_rgb("#1F9B8E"),
        STRIPE=hex_rgb("#1F9B8E"), HAIR=hex_rgb("#E8E8E8"),
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    accents_all = expand_accents(light_accents(), len(usps))
    # Page count is decided ONLY by the count-based split rule -- see the
    # matching comment in render_dark() above.
    pages = chunk_items(list(zip(usps, accents_all)))
    total_pages = len(pages)

    for page_num, page in enumerate(pages, start=1):
        usps_page = [p[0] for p in page]
        accents_page = [p[1] for p in page]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _render_page(slide, args, usps_page, accents_page, page_num, total_pages,
                     theme="light", colors=colors)

    prs.save(out_path_or_prs)
    return prs


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_dir: str) -> list[str]:
    """Convert every slide in pptx_path to a PNG. Returns list of PNG paths
    (one per slide, in order) -- multi-slide-aware (v7), unlike the old
    single-PNG-per-pptx assumption.
    """
    os.makedirs(png_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(pptx_path))[0]
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = sorted(fn for fn in os.listdir(tmp) if fn.startswith("slide-") and fn.endswith(".png"))
        out_paths = []
        for i, fn in enumerate(produced, start=1):
            dest = os.path.join(png_dir, f"{base}_p{i}.png")
            with open(os.path.join(tmp, fn), "rb") as src, open(dest, "wb") as dst:
                dst.write(src.read())
            out_paths.append(dest)
    return out_paths


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True)
    p.add_argument("--theme", required=True, choices=["dark", "light"])
    grp = p.add_mutually_exclusive_group(required=True)
    grp.add_argument("--usps", help="Inline JSON list of {title, description, proof, strategy, enabled} objects (1-5 enabled)")
    grp.add_argument("--usps-json", help="Path to a JSON file containing the USP list")
    p.add_argument("--wedge", default=None,
                   help="Optional manifesto statement, pipe-separated lines (e.g. 'A|B|C').")
    p.add_argument("--wedge-support", default=None,
                   help="Optional one-line supporting sentence under the manifesto.")
    p.add_argument("--out-dir", required=True)
    return p.parse_args()


def main():
    args = parse_args()
    usps = load_usps(args)
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_usp_{args.theme}.pptx")
    if args.theme == "dark":
        render_dark(args, usps, pptx_path)
    else:
        render_light(args, usps, pptx_path)
    pngs = convert_to_png(pptx_path, out_dir)
    print(f"PPTX: {pptx_path}")
    for p in pngs:
        print(f"PNG:  {p}")


if __name__ == "__main__":
    main()
