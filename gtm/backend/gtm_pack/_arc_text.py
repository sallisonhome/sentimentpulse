"""Arc text renderer for the Sizing Rings slide.

python-pptx has no native arc-text primitive, so this module lays out a
string glyph-by-glyph along a circular arc, rotating each glyph so its
baseline is tangent to the arc -- the classic "coin / stamp / badge"
treatment (text hugging the inside of a ring, letters fanning out around
the curve). The result is a fully-editable PPTX (each glyph is a tiny
textbox), not an image overlay -- so themes still apply and the deck
stays editable in PowerPoint if the user wants to tweak.

Geometry / style (v8, badge-stamp pass)
----------------------------------------
Modeled directly on the classic circular-stamp reference: a ring with a
label curving along the TOP half and, if needed, a second label curving
along the BOTTOM half, with small dot separators at 9 o'clock and
3 o'clock between them.

  - TOP arc: text centered at theta=0 (12 o'clock). Letters are set so
    their tops point AWAY from the ring's center (outward) -- reading
    left-to-right along the top of the curve, upright.
  - BOTTOM arc: text centered at theta=180 (6 o'clock). Letters are
    ALSO set so their tops point away from center -- reading
    left-to-right along the bottom of the curve. This is the
    ``flip_for_bottom`` path: mirror image of the top-arc placement,
    achieved by reversing glyph order and adding 180 degrees of
    rotation, so the text is never upside down.
  - Each glyph's rotation exactly equals its own angular position on the
    arc (plus 180 on the bottom), so the glyph's local "up" points
    radially outward at every point along the curve -- true tangency,
    matching the reference image.

Calibration (v9, real-glyph-metrics pass)
------------------------------------------
Earlier versions of this module used a single *average* glyph advance
(``CHAR_ADVANCE_IN_PER_PT = 0.55 / 72.0`` in em terms) for every
character. That is wrong for proportional fonts: wide letters (M, W, N,
O, U, H, G, A, &) take meaningfully more horizontal space than narrow
ones (I, L, J, F, ., ',', space) in Trebuchet MS Bold / Calibri Bold.
Averaging produced visibly uneven letter spacing on the arc -- letters
bunched up around wide characters and stretched out around narrow ones
(e.g. a noticeably bigger gap between "MUST" and "HAVE" than between
"HAVE" and "GAME" in "MUST HAVE GAME").

This version reads REAL per-glyph advance widths from OpenType font
metrics via ``fontTools`` (the ``hmtx`` table), normalized to a
1000-units-per-em basis so widths compose the same way regardless of
the source font's internal ``unitsPerEm``. Widths are cached per
(family, bold) pair so we only touch the filesystem/fontTools once per
font per process.

On a typical Linux rendering host, Trebuchet MS Bold and Calibri Bold
are not installed as native TrueType files (they're Microsoft-licensed
fonts usually only present on Windows/macOS or via `ttf-mscorefonts`).
We use **DejaVu Sans Bold** as the metrics stand-in when available --
it is a widely-installed, metrics-similar sans-serif Bold face that
tracks Trebuchet/Calibri Bold's relative wide/narrow glyph proportions
much better than a flat average does. If DejaVu Sans Bold isn't found
on this system either, we fall back to a hand-calibrated per-character
width table (also normalized to 1000 units/em) so the module still
degrades gracefully rather than crashing.

Coordinate convention
---------------------
All positions and radii are in INCHES. The slide's PPTX coordinate system
has +x to the right and +y DOWN (standard PPTX). Angles are measured with
0 = top, increasing clockwise. So a label above the ring uses theta ~= 0
and a label below the ring uses theta ~= 180.
"""

from __future__ import annotations

import math
import os

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ------------------------------------------------------------------
# Real per-glyph advance widths (fontTools-backed, with a hand-
# calibrated fallback). All widths are normalized to units of 1/1000
# em so downstream math is font-independent: inches = (width_1000 /
# 1000) * pt / 72.
# ------------------------------------------------------------------

# Candidate system font files to use as metrics stand-ins for
# Trebuchet MS Bold / Calibri Bold, in priority order. DejaVu Sans Bold
# is nearly-metrics-similar to both and ships on most Debian/Ubuntu
# systems (used by LibreOffice for headless PPTX->PNG rendering, which
# is exactly the pipeline this module's output flows through).
_FALLBACK_FONT_CANDIDATES = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/msttcorefonts/Trebuchet_MS_Bold.ttf",
    "/usr/share/fonts/truetype/msttcorefonts/Calibri_Bold.ttf",
]

# Hand-calibrated fallback widths (1/1000 em, typical Bold sans-serif
# metrics) -- used only if fontTools/font files are unavailable at all.
CHAR_WIDTHS_EM_1000: dict[str, float] = {
    ' ': 320, '!': 350, '"': 500, '&': 830, "'": 280,
    '(': 400, ')': 400, ',': 320, '-': 400, '.': 320,
    '/': 470, '0': 640, '1': 640, '2': 640, '3': 640,
    '4': 640, '5': 640, '6': 640, '7': 640, '8': 640, '9': 640,
    ':': 350, ';': 350,
    'A': 720, 'B': 720, 'C': 750, 'D': 780, 'E': 660,
    'F': 620, 'G': 800, 'H': 780, 'I': 340, 'J': 460,
    'K': 720, 'L': 620, 'M': 890, 'N': 780, 'O': 800,
    'P': 720, 'Q': 800, 'R': 750, 'S': 680, 'T': 660,
    'U': 780, 'V': 720, 'W': 1000, 'X': 720, 'Y': 700, 'Z': 660,
    # Lowercase (kept in case the caller doesn't uppercase):
    'a': 610, 'b': 630, 'c': 570, 'd': 630, 'e': 600,
    'f': 380, 'g': 630, 'h': 620, 'i': 300, 'j': 340,
    'k': 570, 'l': 300, 'm': 900, 'n': 620, 'o': 620,
    'p': 630, 'q': 630, 'r': 440, 's': 550, 't': 400,
    'u': 620, 'v': 570, 'w': 800, 'x': 570, 'y': 570, 'z': 550,
}
_DEFAULT_CHAR_WIDTH_EM_1000 = 650.0  # unknown glyphs (accents, symbols, etc.)

# Cache: (family, bold) -> dict[str, float] (widths in 1/1000 em)
_FONT_WIDTH_CACHE: dict[tuple[str, bool], dict[str, float]] = {}


def _load_ttfont_widths(path: str) -> dict[str, float] | None:
    """Load per-char advance widths (1/1000 em) from a TTF/OTF file via
    fontTools. Returns None if the file is missing or fontTools can't
    parse it -- callers should fall back to the hand-calibrated table.
    """
    if not path or not os.path.isfile(path):
        return None
    try:
        from fontTools.ttLib import TTFont
    except ImportError:
        return None
    try:
        font = TTFont(path, lazy=True)
        units_per_em = font["head"].unitsPerEm
        hmtx = font["hmtx"]
        cmap = font.getBestCmap()
        scale = 1000.0 / float(units_per_em)
        widths: dict[str, float] = {}
        # Cover the practical character set used on these slides:
        # A-Z, a-z, 0-9, space, and common punctuation.
        chars = (
            "ABCDEFGHIJKLMNOPQRSTUVWXYZ"
            "abcdefghijklmnopqrstuvwxyz"
            "0123456789 !\"&'(),-./:;?"
        )
        for ch in chars:
            cp = ord(ch)
            gname = cmap.get(cp)
            if gname is None:
                continue
            advance = hmtx[gname][0]
            widths[ch] = advance * scale
        return widths if widths else None
    except Exception:
        return None


def _load_font_widths(font_family: str, bold: bool) -> dict[str, float]:
    """Return a dict of per-char advance widths (1/1000 em) for the given
    logical font family. Tries fontTools against known system font
    files first (DejaVu Sans Bold is metrics-close to Trebuchet MS Bold
    / Calibri Bold and is what LibreOffice substitutes on headless
    Linux rendering anyway); falls back to the hand-calibrated table.

    Cached per (family, bold) so repeated calls are cheap.
    """
    key = (font_family, bold)
    cached = _FONT_WIDTH_CACHE.get(key)
    if cached is not None:
        return cached

    widths: dict[str, float] | None = None
    for candidate in _FALLBACK_FONT_CANDIDATES:
        widths = _load_ttfont_widths(candidate)
        if widths:
            break

    if not widths:
        widths = dict(CHAR_WIDTHS_EM_1000)

    _FONT_WIDTH_CACHE[key] = widths
    return widths


def char_width_em1000(ch: str, font_family: str = "Calibri", bold: bool = True) -> float:
    """Advance width of a single character in 1/1000 em units."""
    widths = _load_font_widths(font_family, bold)
    if ch in widths:
        return widths[ch]
    upper = ch.upper()
    if upper in widths:
        return widths[upper]
    return _DEFAULT_CHAR_WIDTH_EM_1000


def char_width_in(ch: str, font_pt: float, font_family: str = "Calibri", bold: bool = True) -> float:
    """Advance width of a single character in INCHES at ``font_pt``."""
    w1000 = char_width_em1000(ch, font_family, bold)
    return (w1000 / 1000.0) * font_pt / 72.0


def text_width_in(text: str, font_pt: float, font_family: str = "Calibri", bold: bool = True) -> float:
    """Total advance width of ``text`` in INCHES at ``font_pt``, summing
    real per-glyph widths rather than an average."""
    if not text:
        return 0.0
    return sum(char_width_in(ch, font_pt, font_family, bold) for ch in text)


def _polar_to_xy(cx: float, cy: float, r: float, theta_deg: float) -> tuple[float, float]:
    """Convert polar (radius, angle-from-12-o'clock-clockwise) to (x,y) inches.

    Angle convention: 0 = up (12 o'clock), 90 = right (3 o'clock),
    180 = down (6 o'clock), 270 = left (9 o'clock). PPTX +y is downward,
    so we FLIP sin's sign in the y calculation.
    """
    a = math.radians(theta_deg)
    x = cx + r * math.sin(a)
    y = cy - r * math.cos(a)
    return x, y


def fit_arc_text_pt(text: str, radius_in: float, max_span_deg: float,
                    base_pt: int = 11, min_pt: int = 6,
                    font: str = "Calibri", bold: bool = True) -> int:
    """Return the largest font size (pt) such that ``text`` fits within
    ``max_span_deg`` when laid out on an arc of ``radius_in`` inches.

    Uses REAL per-glyph advance widths (summed), not an average -- see
    module docstring.
    """
    if not text:
        return base_pt
    max_arc_in = math.radians(max_span_deg) * radius_in * 0.94  # 6% safety
    pt = base_pt
    while pt >= min_pt:
        w_in = text_width_in(text, pt, font, bold)
        if w_in <= max_arc_in:
            return pt
        pt -= 1
    return min_pt


def arc_span_deg(text: str, radius_in: float, pt: float,
                  font: str = "Calibri", bold: bool = True) -> float:
    """Degrees of arc that ``text`` occupies at font size ``pt`` on a ring
    of ``radius_in`` inches."""
    w_in = text_width_in(text, pt, font, bold)
    return math.degrees(w_in / radius_in)


def _place_arc_run(slide, cx: float, cy: float, radius_in: float, text: str,
                    *, theta_mid_deg: float, font: str, pt: float, bold: bool,
                    color: RGBColor | None, is_bottom: bool,
                    letter_spacing_deg_extra: float = 0.0):
    """Place one arc run (already sized) of glyphs centered on theta_mid_deg.

    On the top half (``is_bottom=False``) glyphs are read left-to-right in
    natural order, rotation = own angle (tops point outward/away from
    center, which for the top half means "upright").

    On the bottom half (``is_bottom=True``) glyphs are reversed in render
    order and rotated +180 so that, walking left-to-right along the
    bottom of the ring, the text still reads correctly with letter-tops
    pointing outward (away from center, i.e. downward) -- the standard
    badge/coin mirrored-bottom treatment.

    Each glyph's angular position is computed from the CUMULATIVE
    arc-length of real glyph advances up to that glyph's center, not an
    even/average split -- this is what makes spacing look even across
    strings mixing wide and narrow characters.
    """
    if not text:
        return
    n = len(text)

    # Per-glyph advance widths (inches) in original left-to-right order.
    widths_in = [char_width_in(ch, pt, font, bold) for ch in text]
    total_width_in = sum(widths_in)
    total_span_deg = (math.degrees(total_width_in / radius_in)
                       + letter_spacing_deg_extra * max(0, n - 1))

    # Cumulative center-of-glyph offset (in inches from the start of the
    # string), based on real advances -- glyph i's center sits at
    # (sum of widths before i) + width[i]/2.
    centers_in: list[float] = []
    running = 0.0
    for w in widths_in:
        centers_in.append(running + w / 2.0)
        running += w

    if is_bottom:
        # Reverse render order; mirror the offset axis so the FIRST
        # character in reading order still starts at the correct edge
        # once flipped for bottom-arc placement.
        glyph_order = list(reversed(text))
        offsets_in = [total_width_in - c for c in reversed(centers_in)]
        base_rotation_deg = 180.0
    else:
        glyph_order = list(text)
        offsets_in = centers_in
        base_rotation_deg = 0.0

    if total_width_in <= 0:
        glyph_angles = [theta_mid_deg] * n
    else:
        start_offset_in = total_width_in / 2.0
        glyph_angles = [
            theta_mid_deg + math.degrees((off_in - start_offset_in) / radius_in)
            for off_in in offsets_in
        ]

    glyph_side_in = pt * 1.55 / 72.0  # ample padding, square glyph box

    for glyph, ang in zip(glyph_order, glyph_angles):
        if glyph == " ":
            continue
        gx, gy = _polar_to_xy(cx, cy, radius_in, ang)
        tlx = gx - glyph_side_in / 2.0
        tly = gy - glyph_side_in / 2.0
        box = slide.shapes.add_textbox(
            Inches(tlx), Inches(tly), Inches(glyph_side_in), Inches(glyph_side_in)
        )
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = glyph
        r.font.name = font
        r.font.size = Pt(pt)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
        # Rotation: the glyph's own angular position IS its tangent
        # rotation (0 = upright at 12 o'clock, 90 = rotated 90cw at
        # 3 o'clock, etc). On the bottom arc we add 180 so letters read
        # right-side-up (tops pointing outward/down) rather than upside
        # down.
        rotation_deg = (ang + base_rotation_deg) % 360
        box.rotation = rotation_deg


def add_arc_text(slide, cx: float, cy: float, radius_in: float, text: str,
                 *, theta_mid_deg: float = 0.0, font: str = "Calibri",
                 size: int = 11, bold: bool = True,
                 color: RGBColor | None = None,
                 max_span_deg: float = 160.0,
                 min_pt: int = 6,
                 flip_for_bottom: bool = True):
    """Lay out ``text`` along an arc centered on (cx, cy), stamp/badge style.

    - ``radius_in``: arc radius, in inches -- the visual "text baseline"
      distance from the ring center.
    - ``text``: string to render, uppercase in/out.
    - ``theta_mid_deg``: midpoint of the arc where text is centered. 0 = top,
      180 = bottom.
    - ``max_span_deg``: maximum arc span the text may consume. Font auto-
      shrinks down to ``min_pt`` to fit.
    - ``flip_for_bottom``: when True and the arc midpoint sits on the
      bottom half (90 < theta_mid < 270), the text is placed using the
      mirrored bottom-arc treatment so it reads upright rather than
      upside down.

    Returns the font size actually used.
    """
    if not text:
        return size

    pt = fit_arc_text_pt(text, radius_in, max_span_deg, base_pt=size, min_pt=min_pt,
                          font=font, bold=bold)

    theta_norm = theta_mid_deg % 360
    is_bottom = flip_for_bottom and (90 < theta_norm < 270)

    _place_arc_run(slide, cx, cy, radius_in, text, theta_mid_deg=theta_mid_deg,
                   font=font, pt=pt, bold=bold, color=color, is_bottom=is_bottom)
    return pt


def add_arc_text_split(slide, cx: float, cy: float, radius_in: float, text: str,
                        *, font: str = "Calibri", size: int = 11, bold: bool = True,
                        color: RGBColor | None = None,
                        max_span_deg: float = 200.0, min_pt: int = 6,
                        dot_color: RGBColor | None = None):
    """Badge-stamp layout: try to fit ``text`` on a single TOP arc first.

    If it doesn't fit within ``max_span_deg`` even at ``min_pt``, split the
    words across a TOP arc (theta_mid=0) and a BOTTOM arc (theta_mid=180),
    each spanning up to ~190 degrees, with small dot separators placed at
    9 o'clock (270) and 3 o'clock (90) -- matching the classic circular
    stamp reference where a long label wraps fully around the ring.

    Returns the font size actually used.
    """
    if not text:
        return size

    # Try a single top arc first, at the largest size that still respects
    # min_pt AND keeps total span under max_span_deg.
    single_pt = fit_arc_text_pt(text, radius_in, max_span_deg, base_pt=size, min_pt=min_pt,
                                 font=font, bold=bold)
    single_span = arc_span_deg(text, radius_in, single_pt, font=font, bold=bold)
    if single_span <= max_span_deg or single_pt <= min_pt and single_span <= max_span_deg + 1:
        # Fits on one line around the top -- classic case for short/medium labels.
        if single_span <= max_span_deg:
            add_arc_text(slide, cx, cy, radius_in, text, theta_mid_deg=0.0,
                         font=font, size=size, bold=bold, color=color,
                         max_span_deg=max_span_deg, min_pt=min_pt)
            return single_pt

    # Doesn't fit on a single top arc even at min_pt -- split into two
    # words-balanced halves: top arc + bottom arc, dot separators at the sides.
    words = text.split()
    best_split = len(words) // 2 or 1
    # Choose split point that balances character length between halves.
    if len(words) > 1:
        total = sum(len(w) for w in words) + (len(words) - 1)
        best_diff = None
        for i in range(1, len(words)):
            running = sum(len(w) for w in words[:i]) + (i - 1)
            diff = abs(running - (total - running))
            if best_diff is None or diff < best_diff:
                best_diff = diff
                best_split = i
    top_words = words[:best_split] if len(words) > 1 else words
    bottom_words = words[best_split:] if len(words) > 1 else []
    top_text = " ".join(top_words)
    bottom_text = " ".join(bottom_words)

    half_span = max_span_deg * 0.92  # leave room for dot separators at the sides
    top_pt = fit_arc_text_pt(top_text, radius_in, half_span, base_pt=size, min_pt=min_pt,
                              font=font, bold=bold)
    bottom_pt = (fit_arc_text_pt(bottom_text, radius_in, half_span, base_pt=size, min_pt=min_pt,
                                  font=font, bold=bold)
                 if bottom_text else top_pt)
    used_pt = min(top_pt, bottom_pt)

    _place_arc_run(slide, cx, cy, radius_in, top_text, theta_mid_deg=0.0,
                   font=font, pt=used_pt, bold=bold, color=color, is_bottom=False)
    if bottom_text:
        _place_arc_run(slide, cx, cy, radius_in, bottom_text, theta_mid_deg=180.0,
                       font=font, pt=used_pt, bold=bold, color=color, is_bottom=True)

    # Dot separators at 9 o'clock and 3 o'clock, matching the reference badge.
    if bottom_text:
        dc = dot_color if dot_color is not None else color
        _add_dot(slide, cx, cy, radius_in, 90.0, dc)
        _add_dot(slide, cx, cy, radius_in, 270.0, dc)

    return used_pt


def _add_dot(slide, cx: float, cy: float, radius_in: float, theta_deg: float,
             color: RGBColor | None):
    """Small round separator dot on the arc at the given angle."""
    from pptx.enum.shapes import MSO_SHAPE
    dot_d = 0.045
    gx, gy = _polar_to_xy(cx, cy, radius_in, theta_deg)
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(gx - dot_d / 2), Inches(gy - dot_d / 2),
        Inches(dot_d), Inches(dot_d),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color if color is not None else RGBColor(0, 0, 0)
    s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.text = ""


# ------------------------------------------------------------------
# Center-circle label fitter
# ------------------------------------------------------------------

def _line_width_in(line: str, pt: float, font: str, bold: bool) -> float:
    return text_width_in(line, pt, font, bold)


def _greedy_wrap(text: str, max_width_fn, pt: float, font: str, bold: bool,
                  max_lines: int) -> list[str] | None:
    """Greedy word-wrap ``text`` so each produced line's real rendered
    width (at ``pt``) is <= ``max_width_fn(line_index)`` (a callable so
    each line can have a different usable width based on its
    y-position inside a circle). Returns None if any single word can't
    fit on its own line, or the result needs more than ``max_lines``
    lines.
    """
    words = text.split()
    if not words:
        return []

    lines: list[str] = []
    cur = ""
    line_idx = 0

    def cur_width(candidate: str) -> float:
        return _line_width_in(candidate, pt, font, bold)

    for w in words:
        # A single word wider than the widest possible line can never fit.
        widest_possible = max(max_width_fn(i) for i in range(max_lines))
        if _line_width_in(w, pt, font, bold) > widest_possible:
            return None
        cand = f"{cur} {w}".strip()
        limit = max_width_fn(min(line_idx, max_lines - 1))
        if cur_width(cand) <= limit:
            cur = cand
        else:
            if cur:
                lines.append(cur)
                line_idx += 1
                if line_idx >= max_lines:
                    return None
            cur = w
            # Re-check the single word against its own line's limit.
            limit = max_width_fn(min(line_idx, max_lines - 1))
            if cur_width(cur) > limit:
                return None
    if cur:
        lines.append(cur)

    return lines if len(lines) <= max_lines else None


def fit_inner_circle_pt(text: str, diameter_in: float,
                        base_pt: float = 12, min_pt: float = 5,
                        max_lines: int = 3,
                        font: str = "Calibri", bold: bool = True) -> tuple[float, list[str]]:
    """Return (font_pt, lines) that fit ``text`` fully inside a circle of
    ``diameter_in`` inches, using REAL per-glyph advance widths.

    Strategy: try candidate font sizes from ``base_pt`` down to
    ``min_pt`` in 0.5pt steps. At each size, try wrapping onto 1, then
    2, then 3 lines (up to ``max_lines``). For a given line count N,
    each line's vertical center is estimated from a standard
    single-block-of-text layout centered on the circle's center, and
    the usable chord width for that line is
    ``2 * sqrt(r^2 - y_line^2)`` (shrunk by a small safety margin) --
    the actual width available inside the circle at that height. The
    top and bottom lines must additionally clear the circle's top/
    bottom edge by at least 2% of the diameter. The largest font size
    for which some line-count fits is returned.

    Falls back to a hard character truncation (never overflowing) if
    even ``min_pt`` on ``max_lines`` lines can't fit -- this should be
    extremely rare given the size floor.
    """
    if not text:
        return base_pt, []

    r = diameter_in / 2.0
    top_bottom_pad = 0.02 * diameter_in  # 2% padding rule

    def chord_half_width(y_from_center: float) -> float:
        val = r * r - y_from_center * y_from_center
        return math.sqrt(val) if val > 0 else 0.0

    def usable_width_for_line(y_from_center: float) -> float:
        # 8% inward safety margin so glyphs never visually kiss the
        # circle's curved edge (chord math is exact for the text
        # baseline, but glyph ascenders/descenders and side-bearings
        # need a little breathing room).
        return 2.0 * chord_half_width(y_from_center) * 0.92

    def line_centers(n_lines: int, line_h_in: float) -> list[float]:
        """Vertical center offset (from circle center, +down) for each
        of n_lines lines, stacked and centered as a block."""
        block_h = n_lines * line_h_in
        top_of_block = -block_h / 2.0
        return [top_of_block + (i + 0.5) * line_h_in for i in range(n_lines)]

    def try_size(pt: float) -> list[str] | None:
        line_h_in = (pt / 72.0) * 1.25  # line height incl. natural leading
        for n_lines in range(1, max_lines + 1):
            centers = line_centers(n_lines, line_h_in)
            # Reject candidates where the top or bottom line would cross
            # the circle's top/bottom edge (with the 2% padding rule).
            half_heights = [c + line_h_in / 2.0 for c in centers]
            top_edge = centers[0] - line_h_in / 2.0
            bottom_edge = centers[-1] + line_h_in / 2.0
            if -top_edge > (r - top_bottom_pad):
                continue  # top line pokes above the top edge (with padding)
            if bottom_edge > (r - top_bottom_pad):
                continue  # bottom line pokes below the bottom edge (with padding)

            def max_width_fn(i: int, _centers=centers, _line_h=line_h_in) -> float:
                y = _centers[i]
                # Use the tighter of the two extremes of the line's
                # vertical span, so the whole line height clears the
                # circle boundary, not just its center.
                y_far = max(abs(y - _line_h / 2.0), abs(y + _line_h / 2.0))
                return usable_width_for_line(y_far)

            lines = _greedy_wrap(text, max_width_fn, pt, font, bold, n_lines)
            if lines is not None and len(lines) >= 1:
                return lines
        return None

    pt = base_pt
    while pt >= min_pt:
        lines = try_size(pt)
        if lines is not None:
            return pt, lines
        pt -= 0.5

    # Give up gracefully: min_pt, max_lines, hard character-count wrap
    # sized from the narrowest usable line so we never overflow.
    line_h_in = (min_pt / 72.0) * 1.25
    centers = line_centers(max_lines, line_h_in)
    narrowest = min(usable_width_for_line(abs(c) + line_h_in / 2.0) for c in centers)
    avg_char_w_in = char_width_in("A", min_pt, font, bold)
    max_chars = max(1, int(narrowest / max(avg_char_w_in, 0.001)))
    words = text.split() or [text]
    lines: list[str] = []
    cur = ""
    for w in words:
        cand = f"{cur} {w}".strip()
        if len(cand) <= max_chars:
            cur = cand
        else:
            if cur:
                lines.append(cur)
            cur = w[:max_chars]
        if len(lines) >= max_lines:
            break
    if cur and len(lines) < max_lines:
        lines.append(cur)
    lines = lines[:max_lines]
    return min_pt, lines
