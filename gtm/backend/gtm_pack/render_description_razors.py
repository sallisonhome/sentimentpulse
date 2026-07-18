#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 7 'Game Description & Razors' slide.

Two themes, parity with the rest of the pack:
  - dark   (V2 Modern Mono)  : dark slide, warm gold top bar
  - light  (V4 Bold Brand)   : light slide, left teal stripe

Layout: three cards stacked vertically, each with a heading and a
word-count pill in the top-right corner:
  1. "100-WORD DESCRIPTION" -- prose body (Calibri 11pt, ink)
  2. "20-WORD RAZOR / TAGLINE" -- short tagline (Trebuchet MS Bold 16pt, ink)
  3. "10-WORD RAZOR / SHORT TAGLINE" -- hero line (Trebuchet MS Bold 22pt, accent)

Word counts are informational only (rendered in the pill) -- values over
the nominal limit are WARNED about on stderr but never hard-fail. All
three fields are required (non-empty); that is enforced.

Outputs both a PPTX and a PNG to --out-dir.
"""
from __future__ import annotations

import argparse
import os
import re
import subprocess
import sys
import tempfile

from pptx import Presentation
from ._title_fit import fit_title_pt
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def word_count(text: str) -> int:
    return len(text.split())


def fit_hero_pt(text: str, width_in: float, base_pt: int = 36, min_pt: int = 22) -> int:
    """Return a Trebuchet MS Bold font size (pt) that keeps `text` on ONE line
    within `width_in` inches. Trebuchet Bold at N pt averages ~0.036*N inches
    per character (calibrated empirically against python-pptx output). We shrink
    down toward min_pt in 2pt steps until it fits; if even min_pt overflows we
    return min_pt (caller accepts the wrap).
    """
    if not text:
        return base_pt
    n_chars = len(text)
    pt = base_pt
    while pt >= min_pt:
        est_width = n_chars * 0.036 * pt
        if est_width <= width_in * 0.94:  # 6% safety margin for kerning
            return pt
        pt -= 2
    return min_pt


# ---------- shape primitives ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return box


def add_pill(slide, x, y, text, *, fill, text_color, font="Calibri", size=9):
    """Right-aligned small rounded-rect-look pill (plain rectangle, locked style)."""
    w = 0.35 + 0.085 * len(text)
    pill = add_rect(slide, x - w, y, w, 0.26, fill)
    tf = pill.text_frame
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(9144)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = text_color
    return pill


# ---------- validation ----------
def validate_inputs(args):
    # 2026-07-15: SystemExit is a BaseException, not an Exception, and it's
    # fatal to any FastAPI request that reaches it (the request thread dies
    # before FastAPI's exception handler can wrap it in an HTTPException).
    # Use ValueError so callers can catch cleanly.  CLI usage still surfaces
    # a clear message via argparse's default handling of ValueError below.
    for field, label in (("description", "--description"), ("razor_20", "--razor-20"), ("razor_10", "--razor-10")):
        val = getattr(args, field)
        if not val or not val.strip():
            raise ValueError(f"{label} is required and cannot be empty")

    wc_desc = word_count(args.description)
    wc_r20 = word_count(args.razor_20)
    wc_r10 = word_count(args.razor_10)

    if wc_desc > 100:
        print(f"WARNING: --description is {wc_desc} words (nominal limit 100)", file=sys.stderr)
    if wc_r20 > 20:
        print(f"WARNING: --razor-20 is {wc_r20} words (nominal limit 20)", file=sys.stderr)
    if wc_r10 > 10:
        print(f"WARNING: --razor-10 is {wc_r10} words (nominal limit 10)", file=sys.stderr)

    return wc_desc, wc_r20, wc_r10


# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, counts, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    SURFACE  = hex_rgb("#161A21")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    ACCENT   = hex_rgb("#FFB454")  # warm gold

    wc_desc, wc_r20, wc_r10 = counts

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    add_text(slide, 0.6, 0.4, 10, 0.3, "GAME DESCRIPTION & TAGLINES",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 12, 0.85, args.title,
             font="Trebuchet MS", size=fit_title_pt(args.title, 12), bold=True, color=INK)
    add_text(slide, 0.6, 1.55, 12, 0.4,
             "For product pages, creative briefs, and headlines",
             font="Calibri", size=body_pt(L, 13), color=MUTED)

    cx, cw = 0.6, 12.13

    # Context block -- the 100-word description is de-emphasized (smaller
    # type, no card/border, tightly capped height) since the razors below
    # are the payoff. This is context, not the headline.
    ctx_y = 2.05
    add_text(slide, cx, ctx_y, cw, 0.3, "CONTEXT \u00b7 100-WORD DESCRIPTION",
             font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
    add_text(slide, cx, ctx_y + 0.32, cw - 1.1, 0.85, args.description,
             font="Calibri", size=body_pt(L, 10.5), color=MUTED)
    add_pill(slide, cx + cw, ctx_y, f"{wc_desc} words", fill=SURFACE, text_color=MUTED)

    # Divider between context and the razor payoff
    div_y = ctx_y + 1.28
    add_rect(slide, cx, div_y, cw, 0.012, BORDER)
    add_text(slide, cx, div_y + 0.16, cw, 0.28, "20 & 10 WORD TAGLINES",
             font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # 20-word razor -- secondary payoff, medium size, centered. Word count
    # is a small centered caption below the line (not a floating pill) so
    # it never competes with the centered text for horizontal space.
    # Auto-shrink from 22pt down to 16pt to keep on ONE line.
    r20_y = div_y + 0.62
    r20_pt = fit_hero_pt(args.razor_20, cw, base_pt=22, min_pt=16)
    add_text(slide, cx, r20_y, cw, 0.75, args.razor_20,
             font="Trebuchet MS", size=r20_pt, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, cx, r20_y + 0.62, cw, 0.22, f"{wc_r20} words \u00b7 tagline",
             font="Calibri", size=body_pt(L, 8.5), color=MUTED, align=PP_ALIGN.CENTER)

    # 10-word razor -- the hero line, largest and most dominant element.
    # Auto-shrink to keep it on ONE line -- a wrapped hero line breaks the
    # visual hierarchy that makes this slide work.
    r10_y = r20_y + 1.05
    # Start at 44pt so that even after shrinking to fit a long 10-word razor,
    # the hero remains visibly larger than the 22pt 20-word tagline. Min 26pt
    # (4pt above the 20-word) so hierarchy is preserved even in worst case.
    r10_pt = fit_hero_pt(args.razor_10, cw, base_pt=44, min_pt=26)
    add_text(slide, cx, r10_y, cw, 1.1, args.razor_10,
             font="Trebuchet MS", size=r10_pt, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(slide, cx, r10_y + 1.12, cw, 0.22, f"{wc_r10} words \u00b7 hero tagline",
             font="Calibri", size=body_pt(L, 8.5), color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(out_path)


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, counts, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    PILLBG   = hex_rgb("#F2F2F2")
    C3       = hex_rgb("#1F9B8E")  # structural teal

    wc_desc, wc_r20, wc_r10 = counts

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 0.25, 7.5, C3)

    add_text(slide, 0.7, 0.5, 11, 0.3, "GAME DESCRIPTION & TAGLINES",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3)
    add_text(slide, 0.7, 0.85, 12, 0.85, args.title,
             font="Trebuchet MS", size=fit_title_pt(args.title, 12), bold=True, color=INK)
    add_text(slide, 0.7, 1.65, 12, 0.4,
             "For product pages, creative briefs, and headlines",
             font="Calibri", size=body_pt(L, 14), color=MUTED)

    cx, cw = 0.7, 12.0

    # Context block -- de-emphasized 100-word description
    ctx_y = 2.1
    add_text(slide, cx, ctx_y, cw, 0.3, "CONTEXT \u00b7 100-WORD DESCRIPTION",
             font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
    add_text(slide, cx, ctx_y + 0.32, cw - 1.1, 0.85, args.description,
             font="Calibri", size=body_pt(L, 10.5), color=MUTED)
    add_pill(slide, cx + cw, ctx_y, f"{wc_desc} words", fill=PILLBG, text_color=INK)

    # Divider between context and the razor payoff
    div_y = ctx_y + 1.28
    add_rect(slide, cx, div_y, cw, 0.012, HAIR)
    add_text(slide, cx, div_y + 0.16, cw, 0.28, "20 & 10 WORD TAGLINES",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3, align=PP_ALIGN.CENTER)

    # 20-word razor -- secondary payoff, medium size, centered. Auto-shrink
    # to fit on one line, mirroring the light-theme logic.
    r20_y = div_y + 0.62
    r20_pt = fit_hero_pt(args.razor_20, cw, base_pt=22, min_pt=16)
    add_text(slide, cx, r20_y, cw, 0.75, args.razor_20,
             font="Trebuchet MS", size=r20_pt, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_text(slide, cx, r20_y + 0.62, cw, 0.22, f"{wc_r20} words \u00b7 tagline",
             font="Calibri", size=body_pt(L, 8.5), color=MUTED, align=PP_ALIGN.CENTER)

    # 10-word razor -- the hero line. See render_dark() for the rationale on
    # auto-shrinking; same logic here for the LIGHT theme.
    r10_y = r20_y + 1.05
    r10_pt = fit_hero_pt(args.razor_10, cw, base_pt=44, min_pt=26)
    add_text(slide, cx, r10_y, cw, 1.1, args.razor_10,
             font="Trebuchet MS", size=r10_pt, bold=True, color=C3, align=PP_ALIGN.CENTER)
    add_text(slide, cx, r10_y + 1.12, cw, 0.22, f"{wc_r10} words \u00b7 hero tagline",
             font="Calibri", size=body_pt(L, 8.5), color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(out_path)


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_path: str) -> str:
    out_dir = os.path.dirname(os.path.abspath(png_path))
    os.makedirs(out_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = None
        for fn in sorted(os.listdir(tmp)):
            if fn.startswith("slide-") and fn.endswith(".png"):
                produced = os.path.join(tmp, fn)
                break
        if not produced:
            raise RuntimeError("pdftoppm did not produce a PNG")
        with open(produced, "rb") as src, open(png_path, "wb") as dst:
            dst.write(src.read())
    return png_path


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True, help="Kept for CLI parity; not shown on this slide")
    p.add_argument("--theme", required=True, choices=["dark", "light"])
    p.add_argument("--description", required=True, help="Product-page description, nominally <=100 words")
    p.add_argument("--razor-20", dest="razor_20", required=True, help="Tagline, nominally <=20 words")
    p.add_argument("--razor-10", dest="razor_10", required=True, help="Short hero tagline, nominally <=10 words")
    p.add_argument("--out-dir", required=True)
    return p.parse_args()


def main():
    args = parse_args()
    counts = validate_inputs(args)
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_description_razors_{args.theme}.pptx")
    png_path  = os.path.join(out_dir, f"{slug}_description_razors_{args.theme}.png")
    if args.theme == "dark":
        render_dark(args, counts, pptx_path)
    else:
        render_light(args, counts, pptx_path)
    convert_to_png(pptx_path, png_path)
    print(f"PPTX: {pptx_path}")
    print(f"PNG:  {png_path}")


if __name__ == "__main__":
    main()
