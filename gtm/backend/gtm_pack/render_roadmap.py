#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 4 'Target Activities & Timing' slide series.

GTM Checklist for PC. Generates 5 slides covering all 9 phases (T-12mo+ -> T+365):
  4.1 Foundations & Identity        (T-12+ -> T-9 mo)
  4.2 Store Page & Audience Build   (T-9  -> T-3 mo)
  4.3 Demo Fest & Pre-Launch Push   (T-3  -> T-1 mo)
  4.4 Launch Week & Launch Day      (T-7d -> T+0)
  4.5 Post-Launch Growth            (T+30 -> T+365 days)

Each slide displays:
  - Step header (locked V2/V4 design system)
  - Horizontal phase ribbon at top showing all 5 stages, current stage highlighted,
    with calendar-date anchors computed from --release-date.
  - Two stacked phase-section cards (each with a T-X label, a calendar window,
    and 4-7 checklist items).

Inputs:
  --release-date YYYY-MM-DD  REQUIRED. The target release date. Used to compute
                              the work-back schedule (T-12, T-9, T-6, T-3, T-1,
                              T+0, T+30, T+365).
  --title  --genre           Game title + genre (re-used on every slide).
  --theme  dark | light

Phase content is loaded from assets/roadmap_phases.json (reframed for PC,
no Steam/Valve-specific language).
"""
from __future__ import annotations

import argparse
import datetime as dt
import json
import os
import re
import subprocess
import tempfile
from pathlib import Path

from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt
except ImportError:  # pragma: no cover
    from i18n import body_pt


# ============================================================
# helpers
# ============================================================
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_oval(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
    return box


# ============================================================
# Work-back date math
# ============================================================
def compute_anchors(release_date: dt.date) -> dict:
    """Compute the canonical calendar anchors for the work-back schedule."""
    return {
        "T-12mo": release_date - relativedelta(months=12),
        "T-9mo":  release_date - relativedelta(months=9),
        "T-6mo":  release_date - relativedelta(months=6),
        "T-3mo":  release_date - relativedelta(months=3),
        "T-1mo":  release_date - relativedelta(months=1),
        "T-7d":   release_date - dt.timedelta(days=7),
        "T+0":    release_date,
        "T+30d":  release_date + dt.timedelta(days=30),
        "T+365d": release_date + dt.timedelta(days=365),
    }


def calendar_window(start_label: str, end_label: str, anchors: dict) -> str:
    """Return a human-readable 'Mar 2026 – Jun 2026' window from two anchor keys."""
    def fmt(d: dt.date) -> str:
        return d.strftime("%b %Y")
    s = anchors[start_label]
    e = anchors[end_label]
    if fmt(s) == fmt(e):
        return fmt(s)
    return f"{fmt(s)}  \u2013  {fmt(e)}"


# Mapping each slide -> the two anchor keys that bound its calendar window
SLIDE_WINDOWS = {
    "S4-1": ("T-12mo", "T-9mo"),
    "S4-2": ("T-9mo",  "T-3mo"),
    "S4-3": ("T-3mo",  "T-1mo"),
    "S4-4": ("T-7d",   "T+0"),
    "S4-5": ("T+30d",  "T+365d"),
}

# Mapping each phase section's tx_label -> (start_anchor, end_anchor) for date enrichment
SECTION_WINDOWS = {
    # Slide 1
    "T-12 mo and earlier": ("T-12mo", "T-12mo"),  # single-anchor
    "T-12 \u2192 T-9 mo":  ("T-12mo", "T-9mo"),
    # Slide 2
    "T-9 \u2192 T-6 mo":   ("T-9mo",  "T-6mo"),
    "T-6 \u2192 T-3 mo":   ("T-6mo",  "T-3mo"),
    # Slide 3
    "T-3 mo":              ("T-3mo",  "T-3mo"),
    "T-3 \u2192 T-1 mo":   ("T-3mo",  "T-1mo"),
    # Slide 4
    "T-7 \u2192 T-1 days": ("T-7d",   "T+0"),
    "T+0 (Launch Day)":    ("T+0",    "T+0"),
    # Slide 5
    "T+30 days":           ("T+30d",  "T+30d"),
    "T+60 \u2192 T+365 days": ("T+30d", "T+365d"),
}


def section_date_window(tx_label: str, anchors: dict) -> str:
    if tx_label not in SECTION_WINDOWS:
        return ""
    s, e = SECTION_WINDOWS[tx_label]
    sd, ed = anchors[s], anchors[e]
    fmt = lambda d: d.strftime("%b %Y")
    if tx_label == "T-7 \u2192 T-1 days":
        return f"{sd.strftime('%b %d')} \u2013 {ed.strftime('%b %d, %Y')}"
    if tx_label == "T+0 (Launch Day)":
        return ed.strftime("%a, %b %d, %Y")
    if tx_label == "T+30 days":
        return f"by {ed.strftime('%b %d, %Y')}"
    if tx_label == "T+60 \u2192 T+365 days":
        return f"{(anchors['T+0'] + dt.timedelta(days=60)).strftime('%b %Y')} \u2013 {ed.strftime('%b %Y')}"
    if fmt(sd) == fmt(ed):
        return fmt(sd)
    return f"{fmt(sd)}  \u2013  {fmt(ed)}"


# ============================================================
# Phase ribbon (top of each slide)
# ============================================================
RIBBON_STAGES = [
    ("Foundations",    "T-12+ mo"),
    ("Store & Build",  "T-9 \u2192 T-3"),
    ("Demo & Push",    "T-3 \u2192 T-1"),
    ("Launch",         "T-7d \u2192 T+0"),
    ("Post-Launch",    "T+30 \u2192 T+365"),
]


def draw_ribbon_dark(slide, current_index, anchors, INK, MUTED, BORDER, ACCENT, SURFACE, L="en"):
    """Horizontal ribbon at top showing all 5 stages. Current stage highlighted."""
    rx, ry, rw, rh = 0.6, 2.05, 12.1, 0.55
    add_rect(slide, rx, ry, rw, rh, SURFACE, line=BORDER)
    n = len(RIBBON_STAGES)
    seg_w = rw / n
    # Calendar anchor keys for tick labels under ribbon
    ticks = ["T-12mo", "T-9mo", "T-3mo", "T-1mo", "T+0", "T+365d"]
    for i, (name, txlbl) in enumerate(RIBBON_STAGES):
        x = rx + i * seg_w
        is_curr = (i == current_index)
        # Highlight current segment
        if is_curr:
            add_rect(slide, x + 0.04, ry + 0.06, seg_w - 0.08, rh - 0.12, ACCENT)
            txt_color = hex_rgb("#0E1116")
            sub_color = hex_rgb("#0E1116")
        else:
            txt_color = INK
            sub_color = MUTED
        add_text(slide, x, ry + 0.05, seg_w, 0.28, name,
                 font="Trebuchet MS", size=11, bold=True, color=txt_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ry + 0.30, seg_w, 0.22, txlbl,
                 font="Calibri", size=body_pt(L, 8), color=sub_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # divider between segments
        if i < n - 1:
            add_rect(slide, x + seg_w - 0.005, ry + 0.08, 0.01, rh - 0.16, BORDER)
    # Date tick row under ribbon
    tick_y = ry + rh + 0.05
    label_positions = [0, 0.2, 0.4, 0.6, 0.8, 1.0]
    for keypos, akey in zip(label_positions, ticks):
        tx = rx + keypos * rw - 0.5
        d = anchors[akey].strftime("%b %Y") if "mo" in akey or "+0" in akey else anchors[akey].strftime("%b %Y")
        add_text(slide, tx, tick_y, 1.0, 0.2, d,
                 font="Calibri", size=body_pt(L, 8), color=MUTED, align=PP_ALIGN.CENTER)


def draw_ribbon_light(slide, current_index, anchors, INK, MUTED, HAIR, ACCENT, SURFACE_LIGHT, L="en"):
    rx, ry, rw, rh = 0.7, 2.15, 12.0, 0.55
    add_rect(slide, rx, ry, rw, rh, SURFACE_LIGHT, line=HAIR)
    n = len(RIBBON_STAGES)
    seg_w = rw / n
    ticks = ["T-12mo", "T-9mo", "T-3mo", "T-1mo", "T+0", "T+365d"]
    for i, (name, txlbl) in enumerate(RIBBON_STAGES):
        x = rx + i * seg_w
        is_curr = (i == current_index)
        if is_curr:
            add_rect(slide, x + 0.04, ry + 0.06, seg_w - 0.08, rh - 0.12, ACCENT)
            txt_color = hex_rgb("#FFFFFF")
            sub_color = hex_rgb("#FFFFFF")
        else:
            txt_color = INK
            sub_color = MUTED
        add_text(slide, x, ry + 0.05, seg_w, 0.28, name,
                 font="Trebuchet MS", size=11, bold=True, color=txt_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ry + 0.30, seg_w, 0.22, txlbl,
                 font="Calibri", size=body_pt(L, 8), color=sub_color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_rect(slide, x + seg_w - 0.005, ry + 0.08, 0.01, rh - 0.16, HAIR)
    tick_y = ry + rh + 0.05
    label_positions = [0, 0.2, 0.4, 0.6, 0.8, 1.0]
    for keypos, akey in zip(label_positions, ticks):
        tx = rx + keypos * rw - 0.5
        d = anchors[akey].strftime("%b %Y")
        add_text(slide, tx, tick_y, 1.0, 0.2, d,
                 font="Calibri", size=body_pt(L, 8), color=MUTED, align=PP_ALIGN.CENTER)


# ============================================================
# Slide body: two stacked phase-section cards
# ============================================================
def draw_section_cards_dark(slide, slide_data, anchors, INK, MUTED, BORDER, SURFACE, ACCENT, L="en"):
    """Two phase-section cards, side by side."""
    sections = slide_data["sections"]
    n = len(sections)
    body_y = 3.15
    body_h = 3.8
    margin_x = 0.6
    gap = 0.3
    total_w = 13.333 - 2 * margin_x
    card_w = (total_w - (n - 1) * gap) / n
    for i, sec in enumerate(sections):
        x = margin_x + i * (card_w + gap)
        add_rect(slide, x, body_y, card_w, body_h, SURFACE, line=BORDER)
        # Section heading
        add_text(slide, x + 0.3, body_y + 0.2, card_w - 0.6, 0.4,
                 sec["heading"],
                 font="Trebuchet MS", size=15, bold=True, color=INK)
        # TX label + calendar window on same row
        cal = section_date_window(sec["tx_label"], anchors)
        add_text(slide, x + 0.3, body_y + 0.6, card_w - 0.6, 0.25,
                 sec["tx_label"],
                 font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT)
        if cal:
            add_text(slide, x + 0.3, body_y + 0.85, card_w - 0.6, 0.22,
                     cal, font="Calibri", size=body_pt(L, 9), color=MUTED)
        # Divider
        add_rect(slide, x + 0.3, body_y + 1.13, card_w - 0.6, 0.008, BORDER)
        # Items
        items = sec["items"]
        item_y0 = body_y + 1.25
        item_gap = (body_h - 1.4) / max(len(items), 1)
        for j, item in enumerate(items):
            iy = item_y0 + j * item_gap
            # bullet dot
            add_oval(slide, x + 0.35, iy + 0.13, 0.07, 0.07, ACCENT)
            add_text(slide, x + 0.5, iy, card_w - 0.7, item_gap - 0.05,
                     item, font="Calibri", size=body_pt(L, 10), color=INK)


def draw_section_cards_light(slide, slide_data, anchors, INK, MUTED, HAIR, ACCENT, L="en"):
    sections = slide_data["sections"]
    n = len(sections)
    body_y = 3.25
    body_h = 3.7
    margin_x = 0.7
    gap = 0.3
    total_w = 13.333 - margin_x - 0.6
    card_w = (total_w - (n - 1) * gap) / n
    for i, sec in enumerate(sections):
        x = margin_x + i * (card_w + gap)
        add_rect(slide, x, body_y, card_w, body_h, hex_rgb("#FFFFFF"), line=HAIR)
        # Section heading
        add_text(slide, x + 0.3, body_y + 0.2, card_w - 0.6, 0.4,
                 sec["heading"],
                 font="Trebuchet MS", size=15, bold=True, color=INK)
        cal = section_date_window(sec["tx_label"], anchors)
        add_text(slide, x + 0.3, body_y + 0.6, card_w - 0.6, 0.25,
                 sec["tx_label"],
                 font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT)
        if cal:
            add_text(slide, x + 0.3, body_y + 0.85, card_w - 0.6, 0.22,
                     cal, font="Calibri", size=body_pt(L, 9), color=MUTED)
        add_rect(slide, x + 0.3, body_y + 1.13, card_w - 0.6, 0.008, HAIR)
        items = sec["items"]
        item_y0 = body_y + 1.25
        item_gap = (body_h - 1.4) / max(len(items), 1)
        for j, item in enumerate(items):
            iy = item_y0 + j * item_gap
            add_oval(slide, x + 0.35, iy + 0.13, 0.07, 0.07, ACCENT)
            add_text(slide, x + 0.5, iy, card_w - 0.7, item_gap - 0.05,
                     item, font="Calibri", size=body_pt(L, 10), color=INK)


# ============================================================
# Summary slide (4.6): Key Dates & Events
# ============================================================
# Extra anchor keys used only by the summary slide
EXTRA_ANCHORS = {
    "T-2wk":  lambda r: r - dt.timedelta(days=14),
    "T-1d":   lambda r: r - dt.timedelta(days=1),
    "T+7d":   lambda r: r + dt.timedelta(days=7),
    "T+60d":  lambda r: r + dt.timedelta(days=60),
    "T+90d":  lambda r: r + dt.timedelta(days=90),
    "T+180d": lambda r: r + dt.timedelta(days=180),
}


def resolve_event_date(anchor_key: str, release_date: dt.date, anchors: dict) -> dt.date:
    if anchor_key in anchors:
        return anchors[anchor_key]
    if anchor_key in EXTRA_ANCHORS:
        return EXTRA_ANCHORS[anchor_key](release_date)
    raise ValueError(f"Unknown date anchor: {anchor_key}")


def fmt_event_date(d: dt.date, anchor_key: str) -> str:
    """Format the calendar date for the summary slide. Short, scannable."""
    # Day-level for tight windows; otherwise month + year
    day_level = anchor_key in ("T-2wk", "T-1d", "T+0", "T+7d", "T-7d")
    if day_level:
        return d.strftime("%b %d, %Y")
    return d.strftime("%b %Y")


def draw_summary_dark(slide, events, anchors, args, INK, MUTED, BORDER, SURFACE, ACCENT):
    L = getattr(args, "language", "en")
    """Two-column timeline of date-anchored events. No commentary."""
    # Layout: two columns, events split roughly in half
    body_y = 3.05
    body_h = 4.0
    margin_x = 0.6
    gap = 0.3
    total_w = 13.333 - 2 * margin_x
    card_w = (total_w - gap) / 2

    # Split events: first half pre-launch, second half launch + post-launch
    # Use the index of the first 'T+' event as the natural pivot
    pivot = next((i for i, e in enumerate(events) if e["date_anchor"].startswith("T+")), len(events) // 2)
    columns = [events[:pivot], events[pivot:]]
    col_titles = ["PRE-LAUNCH", "LAUNCH & POST-LAUNCH"]

    for col_i, col_events in enumerate(columns):
        x = margin_x + col_i * (card_w + gap)
        add_rect(slide, x, body_y, card_w, body_h, SURFACE, line=BORDER)
        # Column title
        add_text(slide, x + 0.3, body_y + 0.2, card_w - 0.6, 0.3,
                 col_titles[col_i],
                 font="Trebuchet MS", size=11, bold=True, color=ACCENT)
        # Divider
        add_rect(slide, x + 0.3, body_y + 0.58, card_w - 0.6, 0.008, BORDER)

        # Date column width within card
        date_col_w = 1.55
        row_y0 = body_y + 0.75
        row_gap = (body_h - 0.95) / max(len(col_events), 1)
        for j, ev in enumerate(col_events):
            ry = row_y0 + j * row_gap
            d = resolve_event_date(ev["date_anchor"], args.release_date, anchors)
            date_str = fmt_event_date(d, ev["date_anchor"])
            # Date pill (left)
            add_text(slide, x + 0.3, ry, date_col_w, row_gap - 0.02,
                     date_str, font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT,
                     anchor=MSO_ANCHOR.MIDDLE)
            # Event label (right)
            add_text(slide, x + 0.3 + date_col_w + 0.1, ry,
                     card_w - 0.6 - date_col_w - 0.1, row_gap - 0.02,
                     ev["event"], font="Calibri", size=body_pt(L, 10), color=INK,
                     anchor=MSO_ANCHOR.MIDDLE)


def draw_summary_light(slide, events, anchors, args, INK, MUTED, HAIR, ACCENT):
    L = getattr(args, "language", "en")
    body_y = 3.15
    body_h = 3.9
    margin_x = 0.7
    gap = 0.3
    total_w = 13.333 - margin_x - 0.6
    card_w = (total_w - gap) / 2

    pivot = next((i for i, e in enumerate(events) if e["date_anchor"].startswith("T+")), len(events) // 2)
    columns = [events[:pivot], events[pivot:]]
    col_titles = ["PRE-LAUNCH", "LAUNCH & POST-LAUNCH"]

    for col_i, col_events in enumerate(columns):
        x = margin_x + col_i * (card_w + gap)
        add_rect(slide, x, body_y, card_w, body_h, hex_rgb("#FFFFFF"), line=HAIR)
        add_text(slide, x + 0.3, body_y + 0.2, card_w - 0.6, 0.3,
                 col_titles[col_i],
                 font="Trebuchet MS", size=11, bold=True, color=ACCENT)
        add_rect(slide, x + 0.3, body_y + 0.58, card_w - 0.6, 0.008, HAIR)

        date_col_w = 1.55
        row_y0 = body_y + 0.75
        row_gap = (body_h - 0.95) / max(len(col_events), 1)
        for j, ev in enumerate(col_events):
            ry = row_y0 + j * row_gap
            d = resolve_event_date(ev["date_anchor"], args.release_date, anchors)
            date_str = fmt_event_date(d, ev["date_anchor"])
            add_text(slide, x + 0.3, ry, date_col_w, row_gap - 0.02,
                     date_str, font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT,
                     anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, x + 0.3 + date_col_w + 0.1, ry,
                     card_w - 0.6 - date_col_w - 0.1, row_gap - 0.02,
                     ev["event"], font="Calibri", size=body_pt(L, 10), color=INK,
                     anchor=MSO_ANCHOR.MIDDLE)


def draw_summary_ribbon_dark(slide, anchors, INK, MUTED, BORDER, ACCENT, SURFACE, L="en"):
    """Same ribbon as content slides, but no single phase highlighted (overview)."""
    rx, ry, rw, rh = 0.6, 2.05, 12.1, 0.55
    add_rect(slide, rx, ry, rw, rh, SURFACE, line=BORDER)
    n = len(RIBBON_STAGES)
    seg_w = rw / n
    # Full accent underline to signal this slide spans all phases
    add_rect(slide, rx, ry + rh - 0.06, rw, 0.06, ACCENT)
    ticks = ["T-12mo", "T-9mo", "T-3mo", "T-1mo", "T+0", "T+365d"]
    for i, (name, txlbl) in enumerate(RIBBON_STAGES):
        x = rx + i * seg_w
        add_text(slide, x, ry + 0.05, seg_w, 0.28, name,
                 font="Trebuchet MS", size=11, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ry + 0.30, seg_w, 0.22, txlbl,
                 font="Calibri", size=body_pt(L, 8), color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_rect(slide, x + seg_w - 0.005, ry + 0.08, 0.01, rh - 0.16, BORDER)
    tick_y = ry + rh + 0.05
    label_positions = [0, 0.2, 0.4, 0.6, 0.8, 1.0]
    for keypos, akey in zip(label_positions, ticks):
        tx = rx + keypos * rw - 0.5
        d = anchors[akey].strftime("%b %Y")
        add_text(slide, tx, tick_y, 1.0, 0.2, d,
                 font="Calibri", size=body_pt(L, 8), color=MUTED, align=PP_ALIGN.CENTER)


def draw_summary_ribbon_light(slide, anchors, INK, MUTED, HAIR, ACCENT, SURFACE_LIGHT, L="en"):
    rx, ry, rw, rh = 0.7, 2.15, 12.0, 0.55
    add_rect(slide, rx, ry, rw, rh, SURFACE_LIGHT, line=HAIR)
    n = len(RIBBON_STAGES)
    seg_w = rw / n
    add_rect(slide, rx, ry + rh - 0.06, rw, 0.06, ACCENT)
    ticks = ["T-12mo", "T-9mo", "T-3mo", "T-1mo", "T+0", "T+365d"]
    for i, (name, txlbl) in enumerate(RIBBON_STAGES):
        x = rx + i * seg_w
        add_text(slide, x, ry + 0.05, seg_w, 0.28, name,
                 font="Trebuchet MS", size=11, bold=True, color=INK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, ry + 0.30, seg_w, 0.22, txlbl,
                 font="Calibri", size=body_pt(L, 8), color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_rect(slide, x + seg_w - 0.005, ry + 0.08, 0.01, rh - 0.16, HAIR)
    tick_y = ry + rh + 0.05
    label_positions = [0, 0.2, 0.4, 0.6, 0.8, 1.0]
    for keypos, akey in zip(label_positions, ticks):
        tx = rx + keypos * rw - 0.5
        d = anchors[akey].strftime("%b %Y")
        add_text(slide, tx, tick_y, 1.0, 0.2, d,
                 font="Calibri", size=body_pt(L, 8), color=MUTED, align=PP_ALIGN.CENTER)


def render_summary_dark(prs, events, slide_num, total, anchors, args):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    SURFACE  = hex_rgb("#161A21")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    ACCENT   = hex_rgb("#FFB454")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    add_text(slide, 0.6, 0.4, 12, 0.3,
             f"STEP 04 \u00b7 TARGET ACTIVITIES & TIMING  \u00b7  {slide_num}/{total}",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 12, 0.6, f"Key Dates & Events \u2014 {args.title}",
             font="Trebuchet MS", size=28, bold=True, color=INK)
    add_text(slide, 0.6, 1.35, 12, 0.4,
             f"{args.genre}  \u00b7  Work-back summary. Calendar-anchored milestones only.",
             font="Calibri", size=body_pt(L, 12), color=MUTED)
    rel_str = "Release: " + args.release_date.strftime("%b %Y")
    add_rect(slide, 10.4, 0.4, 2.3, 0.32, SURFACE, line=BORDER)
    add_text(slide, 10.4, 0.4, 2.3, 0.32, rel_str,
             font="Calibri", size=body_pt(L, 9), bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    draw_summary_ribbon_dark(slide, anchors, INK, MUTED, BORDER, ACCENT, SURFACE, L)
    draw_summary_dark(slide, events, anchors, args, INK, MUTED, BORDER, SURFACE, ACCENT)

    add_text(slide, 0.6, 7.15, 12, 0.25,
             f"GTM CHECKLIST FOR PC  \u00b7  STEP 04.{slide_num}  \u00b7  WORK-BACK FROM {args.release_date.strftime('%b %d, %Y').upper()}",
             font="Calibri", size=body_pt(L, 8), bold=True, color=MUTED)


def render_summary_light(prs, events, slide_num, total, anchors, args):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    SURFACE_LIGHT = hex_rgb("#FAFAF7")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    ACCENT   = hex_rgb("#1F9B8E")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 0.25, 7.5, ACCENT)

    add_text(slide, 0.7, 0.5, 12, 0.3,
             f"STEP 04 \u00b7 TARGET ACTIVITIES & TIMING  \u00b7  {slide_num}/{total}",
             font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT)
    add_text(slide, 0.7, 0.85, 12, 0.6, f"Key Dates & Events \u2014 {args.title}",
             font="Trebuchet MS", size=28, bold=True, color=INK)
    add_text(slide, 0.7, 1.45, 12, 0.4,
             f"{args.genre} \u00b7 Work-back summary. Calendar-anchored milestones only.",
             font="Calibri", size=body_pt(L, 12), color=MUTED)
    rel_str = "Release: " + args.release_date.strftime("%b %Y")
    add_rect(slide, 10.4, 0.5, 2.3, 0.32, SURFACE_LIGHT, line=HAIR)
    add_text(slide, 10.4, 0.5, 2.3, 0.32, rel_str,
             font="Calibri", size=body_pt(L, 9), bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    draw_summary_ribbon_light(slide, anchors, INK, MUTED, HAIR, ACCENT, SURFACE_LIGHT, L)
    draw_summary_light(slide, events, anchors, args, INK, MUTED, HAIR, ACCENT)

    add_text(slide, 0.7, 7.15, 12, 0.25,
             f"GTM CHECKLIST FOR PC  \u00b7  STEP 04.{slide_num}  \u00b7  WORK-BACK FROM {args.release_date.strftime('%b %d, %Y').upper()}",
             font="Calibri", size=body_pt(L, 8), bold=True, color=MUTED)


# ============================================================
# Render one slide (dark)
# ============================================================
def render_one_dark(prs, slide_data, slide_num, total, anchors, args):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    SURFACE  = hex_rgb("#161A21")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    ACCENT   = hex_rgb("#FFB454")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background + top accent bar
    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    # Header
    add_text(slide, 0.6, 0.4, 12, 0.3,
             f"STEP 04 \u00b7 TARGET ACTIVITIES & TIMING  \u00b7  {slide_num}/{total}",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 12, 0.6, f"{slide_data['title']} \u2014 {args.title}",
             font="Trebuchet MS", size=28, bold=True, color=INK)
    add_text(slide, 0.6, 1.35, 12, 0.4,
             f"{args.genre}  \u00b7  {slide_data['subtitle']}",
             font="Calibri", size=body_pt(L, 12), color=MUTED)
    # Release date pill (top-right)
    rel_str = "Release: " + args.release_date.strftime("%b %Y")
    add_rect(slide, 10.4, 0.4, 2.3, 0.32, SURFACE, line=BORDER)
    add_text(slide, 10.4, 0.4, 2.3, 0.32, rel_str,
             font="Calibri", size=body_pt(L, 9), bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Phase ribbon
    draw_ribbon_dark(slide, slide_num - 1, anchors, INK, MUTED, BORDER, ACCENT, SURFACE, L)

    # Body cards
    draw_section_cards_dark(slide, slide_data, anchors, INK, MUTED, BORDER, SURFACE, ACCENT, L)

    # Footer
    add_text(slide, 0.6, 7.15, 12, 0.25,
             f"GTM CHECKLIST FOR PC  \u00b7  STEP 04.{slide_num}  \u00b7  WORK-BACK FROM {args.release_date.strftime('%b %d, %Y').upper()}",
             font="Calibri", size=body_pt(L, 8), bold=True, color=MUTED)


# ============================================================
# Render one slide (light)
# ============================================================
def render_one_light(prs, slide_data, slide_num, total, anchors, args):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    SURFACE_LIGHT = hex_rgb("#FAFAF7")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    ACCENT   = hex_rgb("#1F9B8E")

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    # Left accent stripe (locked light motif)
    add_rect(slide, 0, 0, 0.25, 7.5, ACCENT)

    # Header
    add_text(slide, 0.7, 0.5, 12, 0.3,
             f"STEP 04 \u00b7 TARGET ACTIVITIES & TIMING  \u00b7  {slide_num}/{total}",
             font="Calibri", size=body_pt(L, 10), bold=True, color=ACCENT)
    add_text(slide, 0.7, 0.85, 12, 0.6, f"{slide_data['title']} \u2014 {args.title}",
             font="Trebuchet MS", size=28, bold=True, color=INK)
    add_text(slide, 0.7, 1.45, 12, 0.4,
             f"{args.genre} \u00b7 {slide_data['subtitle']}",
             font="Calibri", size=body_pt(L, 12), color=MUTED)
    # Release pill (top-right)
    rel_str = "Release: " + args.release_date.strftime("%b %Y")
    add_rect(slide, 10.4, 0.5, 2.3, 0.32, SURFACE_LIGHT, line=HAIR)
    add_text(slide, 10.4, 0.5, 2.3, 0.32, rel_str,
             font="Calibri", size=body_pt(L, 9), bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    draw_ribbon_light(slide, slide_num - 1, anchors, INK, MUTED, HAIR, ACCENT, SURFACE_LIGHT, L)
    draw_section_cards_light(slide, slide_data, anchors, INK, MUTED, HAIR, ACCENT, L)

    add_text(slide, 0.7, 7.15, 12, 0.25,
             f"GTM CHECKLIST FOR PC  \u00b7  STEP 04.{slide_num}  \u00b7  WORK-BACK FROM {args.release_date.strftime('%b %d, %Y').upper()}",
             font="Calibri", size=body_pt(L, 8), bold=True, color=MUTED)


# ============================================================
# Main
# ============================================================
def render_pptx_to_png(pptx_path, out_dir):
    """Convert each slide of pptx to PNG via LibreOffice + pdftoppm."""
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                        "--outdir", tmp, pptx_path], check=True, capture_output=True)
        base = Path(pptx_path).stem
        pdf_path = os.path.join(tmp, f"{base}.pdf")
        prefix = os.path.join(out_dir, base)
        subprocess.run(["pdftoppm", "-png", "-r", "150", pdf_path, prefix],
                       check=True, capture_output=True)


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True, help="Game title")
    p.add_argument("--genre", required=True, help="Genre tagline")
    p.add_argument("--release-date", required=True,
                   help="Target release date YYYY-MM-DD (work-back anchor)")
    p.add_argument("--theme", choices=["dark", "light"], required=True)
    p.add_argument("--phases-json", default=None,
                   help="Override path to roadmap_phases.json")
    p.add_argument("--out-dir", required=True)
    p.add_argument("--no-png", action="store_true",
                   help="Skip PNG rendering (PPTX only)")
    args = p.parse_args()

    args.release_date = dt.date.fromisoformat(args.release_date)

    # Locate phases JSON
    if args.phases_json:
        phases_path = args.phases_json
    else:
        here = Path(__file__).resolve().parent.parent
        phases_path = here / "assets" / "roadmap_phases.json"
    with open(phases_path, "r") as f:
        phases = json.load(f)

    anchors = compute_anchors(args.release_date)

    os.makedirs(args.out_dir, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = phases["slides"]
    summary_events = phases.get("summary_events", {}).get("events", [])
    total = len(slides_data) + (1 if summary_events else 0)
    for i, sd in enumerate(slides_data):
        if args.theme == "dark":
            render_one_dark(prs, sd, i + 1, total, anchors, args)
        else:
            render_one_light(prs, sd, i + 1, total, anchors, args)
    # Summary slide (4.6)
    if summary_events:
        sn = len(slides_data) + 1
        if args.theme == "dark":
            render_summary_dark(prs, summary_events, sn, total, anchors, args)
        else:
            render_summary_light(prs, summary_events, sn, total, anchors, args)

    title_slug = slugify(args.title)
    pptx_name = f"step4_roadmap_{title_slug}_{args.theme}.pptx"
    pptx_path = os.path.join(args.out_dir, pptx_name)
    prs.save(pptx_path)
    print(f"Wrote: {pptx_path}")

    if not args.no_png:
        render_pptx_to_png(pptx_path, args.out_dir)
        print(f"PNGs written to: {args.out_dir}")


if __name__ == "__main__":
    main()
