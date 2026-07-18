#!/usr/bin/env python3
"""Render the GTM Slide Pack Step 1 'Target Audiences & Sizing' nested-circle slide.

Two themes baked in:
  - dark   (V2 Modern Mono)  : dark slide, teal ramp + warm gold accent on breakout tier, KPI cards
  - light  (V4 Bold Brand)   : light slide, refined original palette, color-swatch legend rows

v7 polish pass (2026-07-18): every ring is now filled with its tier color and
carries its own name + formatted count directly inside (or, for the outer
ring where the band is too thin, just outside the ring on a matching-color
chip) with a contrast-aware foreground color. The side KPI/legend card is
KEPT next to the chart. Subtitle now explicitly calls out "Potential Buyers"
so the numbers are self-explanatory. Footer removed (global v7 rule).

Outputs both a PPTX and a PNG to --out-dir.
"""
from __future__ import annotations

import argparse
import os
import re
import subprocess
import tempfile

from pptx import Presentation
from ._title_fit import fit_title_pt, wrap_label
from ._arc_text import add_arc_text, add_arc_text_split, arc_span_deg, fit_arc_text_pt, fit_inner_circle_pt
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

try:
    from .i18n import body_pt  # package-relative (used via gtm_pack/__init__.py wrapper)
except ImportError:  # pragma: no cover - direct-script invocation fallback
    from i18n import body_pt


# ---------- helpers ----------
def hex_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def relative_luminance(h: str) -> float:
    """Approximate relative luminance (0=black, 1=white) for contrast picks."""
    h = h.lstrip("#")
    r, g, b = int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def fg_for_bg(hex_bg: str, light_fg: str, dark_fg: str) -> RGBColor:
    """Pick light_fg or dark_fg (hex strings) based on background luminance."""
    return hex_rgb(light_fg) if relative_luminance(hex_bg) < 0.55 else hex_rgb(dark_fg)


def slugify(value: str) -> str:
    s = re.sub(r"[^A-Za-z0-9]+", "_", value.strip().lower()).strip("_")
    return s or "untitled"


def fmt_num(n: int) -> str:
    return f"{n:,}"


def fmt_short(n: int) -> str:
    """Format a number compactly for in-circle labels (1.2M, 850K, 12.3M)."""
    if n >= 1_000_000:
        v = n / 1_000_000
        return f"{v:.1f}M" if v < 10 else f"{int(round(v))}M"
    if n >= 1_000:
        v = n / 1_000
        return f"{int(round(v))}K"
    return f"{n:,}"


# ---------- shape primitives ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_circle(slide, cx, cy, d, fill, line=None, line_w_pt=2.0):
    s = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(cx - d / 2), Inches(cy - d / 2), Inches(d), Inches(d)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w_pt)
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def add_text(slide, x, y, w, h, text, *, font="Calibri", size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
             line_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color
    return box


# ---------- Cohort framing ----------
def innermost_label(args) -> tuple[str, str]:
    """Return (chart_label, legend_label) for the innermost cohort."""
    if args.inner == "prev":
        return ("Prev owners", "Prev Game Owners")
    if args.inner == "dev":
        return ("Dev fans", "Developer Fans")
    # custom
    name = args.inner_name or "Cohort"
    return (name, name)


def ring2_chart_label(args) -> str:
    if args.type == "custom":
        return args.ring2_name or "Custom cohort"
    return "IP Fans"


def ring2_legend_name(args) -> str:
    if args.type == "custom":
        return args.ring2_name
    return "IP Fans (no prior)"


def ring2_legend_desc(args) -> str:
    if args.type == "custom":
        return args.ring2_definition or ""
    return "Followers of the IP who didn't own previous"


def cohort3_name(args) -> str:
    """User-typed name for ring 3 (was hardcoded 'Genre Fans'). Falls back
    to the 'Genre Fans' default ONLY when no user-typed name was supplied --
    the wizard's Step 3 always collects a name for this cohort, so this
    fallback should rarely trigger in practice."""
    return getattr(args, "cohort3_name", None) or "Genre Fans"


def cohort3_desc(args) -> str:
    return getattr(args, "cohort3_definition", None) or f"Top 5 avg in {args.genre}"


def cohort4_name(args) -> str:
    """User-typed name for the outer ring (was hardcoded 'Breakout Ceiling')."""
    return getattr(args, "cohort4_name", None) or "Breakout Ceiling"


def cohort4_desc(args) -> str:
    return getattr(args, "cohort4_definition", None) or f"Top 2 ever in {args.genre}+"


# ============================================================
# Shared ring arc-label renderer (used by BOTH themes)
# ============================================================
def render_ring_label(slide, cx, cy, r_outer, r_inner, label_text,
                       *, fg_color, base_pt=11, min_pt=7,
                       label_font="Calibri", dot_color=None):
    """Render one ring's badge-stamp arc label -- poster/coin style.

    Per the locked design decision, the ring graphic carries ONLY the
    curved cohort-name label; the numeric counts live exclusively in the
    side card ("POTENTIAL BUYERS BY TIER"), so there is no number to make
    room for here. That frees up the full ring band for a bigger, bolder
    arc label than earlier iterations.

    - The arc label curves along the TOP of the ring band, following the
      classic circular-stamp reference: baseline tangent to the arc,
      letters upright at 12 o'clock, rotated progressively toward the
      sides.
    - If the label is too long to fit on a single top arc even at
      ``min_pt``, it automatically splits into a TOP arc + BOTTOM arc
      with dot separators at 9 and 3 o'clock (still stamp-style) -- the
      label wraps fully around the ring rather than truncating.

    Returns the font size used for the label.
    """
    band_width = r_outer - r_inner
    # Label radius sits at the visual center of the band (not hugging the
    # rim) now that the band doesn't have to also host a number -- this
    # lets the label read as a bigger, bolder badge inscription.
    label_radius = r_outer - band_width * 0.42
    # Max span before we consider this ring "full" -- leaves a small gap
    # at the very sides of the ring (badge labels never wrap fully to 360).
    max_span_deg = 214.0

    text_upper = label_text.upper()
    single_pt = fit_arc_text_pt(text_upper, label_radius, max_span_deg,
                                 base_pt=base_pt, min_pt=min_pt)
    single_span = arc_span_deg(text_upper, label_radius, single_pt)

    if single_span <= max_span_deg:
        used_pt = add_arc_text(slide, cx, cy, label_radius, text_upper,
                                theta_mid_deg=0.0, font=label_font,
                                size=base_pt, bold=True, color=fg_color,
                                max_span_deg=max_span_deg, min_pt=min_pt)
    else:
        # Split top/bottom -- the label wraps the FULL ring (top + bottom
        # arcs with dot separators at 9 and 3 o'clock).
        used_pt = add_arc_text_split(slide, cx, cy, label_radius, text_upper,
                                      font=label_font, size=base_pt, bold=True,
                                      color=fg_color, max_span_deg=max_span_deg,
                                      min_pt=min_pt,
                                      dot_color=dot_color or fg_color)
    return used_pt


def render_center_label(slide, cx, cy, diameter_in, label_text, *, fg_color,
                         font="Calibri", base_pt=20, min_pt=10):
    """Render ONLY the cohort-name label inside the center circle --
    bigger and vertically centered now that there's no number to share
    room with. Auto-shrinks and wraps to at most 2 lines so it can never
    break outside the circle's bounds, however long the user's cohort
    name is.
    """
    text_upper = label_text.upper()
    label_pt, label_lines = fit_inner_circle_pt(text_upper, diameter_in,
                                                 base_pt=base_pt, min_pt=min_pt,
                                                 max_lines=2)
    if not label_lines:
        return label_pt
    lbl_h = (label_pt / 72.0) * 1.35 * len(label_lines) + 0.06
    add_text(slide, cx - diameter_in * 0.46, cy - lbl_h / 2, diameter_in * 0.92, lbl_h,
             label_lines,
             font=font, size=label_pt, bold=True, color=fg_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return label_pt


# ============================================================
# THEME: DARK  (V2 Modern Mono)
# ============================================================
def render_dark(args, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#0E1116")
    SURFACE  = hex_rgb("#161A21")
    BORDER   = hex_rgb("#1F2530")
    INK      = hex_rgb("#E8E6E1")
    MUTED    = hex_rgb("#8A8F99")
    A1       = hex_rgb("#0A2A30")  # outer
    A2       = hex_rgb("#155966")
    A3       = hex_rgb("#2FA9BD")
    A4       = hex_rgb("#7FD8E3")  # inner
    ACCENT   = hex_rgb("#FFB454")  # warm gold — used on breakout

    A1_HEX, A2_HEX, A3_HEX, A4_HEX = "#0A2A30", "#155966", "#2FA9BD", "#7FD8E3"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background + top accent bar
    add_rect(slide, 0, 0, 13.333, 7.5, BG)
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT)

    # Title block
    add_text(slide, 0.6, 0.4, 8, 0.3, "TARGET AUDIENCES & SIZING",
             font="Trebuchet MS", size=10, bold=True, color=ACCENT)
    add_text(slide, 0.6, 0.75, 11, 0.85, args.title,
             font="Trebuchet MS", size=fit_title_pt(args.title, 11), bold=True, color=INK)
    type_label = {"sequel": "Sequel", "new_ip_with_fans": "IP-based",
                  "custom": "Original IP"}[args.type]
    add_text(slide, 0.6, 1.55, 11.5, 0.4,
             f"Potential buyer audience by tier · {args.genre}  ·  {type_label}",
             font="Calibri", size=body_pt(L, 13), color=MUTED)

    # ---- Circles (left half) — filled, labeled tiers ----
    cx, cy = 3.55, 4.75
    d1, d2, d3, d4 = 5.2, 3.95, 2.65, 1.35  # outer -> inner diameters
    add_circle(slide, cx, cy, d1, A1, line=BORDER, line_w_pt=0.75)
    add_circle(slide, cx, cy, d2, A2, line=BORDER, line_w_pt=0.75)
    add_circle(slide, cx, cy, d3, A3, line=BORDER, line_w_pt=0.75)
    add_circle(slide, cx, cy, d4, A4, line=BORDER, line_w_pt=0.75)

    inner_chart = ring2_chart_label(args)
    prev_chart, _ = innermost_label(args)

    fg1 = fg_for_bg(A1_HEX, "#E8E6E1", "#0E1116")
    fg2 = fg_for_bg(A2_HEX, "#E8E6E1", "#0E1116")
    fg3 = fg_for_bg(A3_HEX, "#E8E6E1", "#0E1116")
    fg4 = fg_for_bg(A4_HEX, "#E8E6E1", "#0E1116")

    # Ring radii, outer -> inner.
    r1, r2, r3, r4 = d1 / 2, d2 / 2, d3 / 2, d4 / 2

    # v9 pass: numeric counts live ONLY in the side card now (removed from
    # the ring graphic per locked design decision) -- each of the three
    # outer rings gets ONLY a curved arc-text label (classic circular-
    # stamp treatment, see render_ring_label / _arc_text.py), sized bigger
    # now that the full band is available for the label alone. Long labels
    # that don't fit a single top arc automatically split into top+bottom
    # arcs with dot separators, so nothing ever truncates.

    # ---- Ring 1 (outermost, breakout) ----
    c4_label = cohort4_name(args)
    render_ring_label(slide, cx, cy, r1, r2, c4_label,
                       fg_color=fg1, base_pt=15, min_pt=9)

    # ---- Ring 2 (cohort 3 / "genre fans" slot) ----
    c3_label = cohort3_name(args)
    render_ring_label(slide, cx, cy, r2, r3, c3_label,
                       fg_color=fg2, base_pt=13, min_pt=8)

    # ---- Ring 3 (cohort 2 / ip fans / custom ring2) ----
    render_ring_label(slide, cx, cy, r3, r4, inner_chart,
                       fg_color=fg3, base_pt=11.5, min_pt=7)

    # ---- Inner circle (cohort 1) ----
    # v9: number removed (lives only in the side card now) -- the center
    # circle shows ONLY the cohort-name label, bigger and centered, still
    # auto-shrinking/wrapping so it NEVER breaks outside the circle bounds.
    render_center_label(slide, cx, cy, d4, prev_chart, fg_color=fg4)

    # ---- KPI cards (right half) — kept, per user's "keep side card" pick ----
    _, inner_legend = innermost_label(args)
    inner_desc = "Players of the prior title" if args.inner == "prev" else \
                 "Direct followers of the developer" if args.inner == "dev" else \
                 (args.inner_definition or "")

    cards = [
        (inner_legend.upper(),               fmt_num(args.prev),       inner_desc,                              A4),
        (ring2_legend_name(args).upper(),    fmt_num(args.ip_fans),    ring2_legend_desc(args),                 A3),
        (cohort3_name(args).upper(),          fmt_num(args.genre_fans), cohort3_desc(args),                      A2),
        (cohort4_name(args).upper(),          fmt_num(args.breakout),   cohort4_desc(args),                      ACCENT),
    ]
    rx = 7.5
    ry = 2.15
    cw = 5.25
    ch = 1.08
    gap = 0.14
    add_text(slide, rx, ry - 0.32, cw, 0.28, "POTENTIAL BUYERS BY TIER",
             font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
    for i, (label, num, desc, accent) in enumerate(cards):
        y = ry + i * (ch + gap)
        add_rect(slide, rx, y, cw, ch, SURFACE)
        # Tier accent strip — left side, indicates which ring this row maps to
        add_rect(slide, rx, y, 0.06, ch, accent)
        add_text(slide, rx + 0.3, y + 0.11, cw - 0.5, 0.28, label,
                 font="Calibri", size=body_pt(L, 9), bold=True, color=MUTED)
        add_text(slide, rx + 0.3, y + 0.38, cw - 0.5, 0.42, num,
                 font="Trebuchet MS", size=20, bold=True, color=INK)
        add_text(slide, rx + 0.3, y + 0.78, cw - 0.5, 0.25, desc,
                 font="Calibri", size=body_pt(L, 9.5), color=MUTED)

    prs.save(out_path)


# ============================================================
# THEME: LIGHT (V4 Bold Brand)
# ============================================================
def render_light(args, out_path):
    L = getattr(args, "language", "en")
    BG       = hex_rgb("#FFFFFF")
    INK      = hex_rgb("#1A1A1A")
    MUTED    = hex_rgb("#5C5C5C")
    HAIR     = hex_rgb("#E8E8E8")
    # Refined palette — same hue family as user's original template
    C1       = hex_rgb("#E5A700")  # gold (outer)
    C2       = hex_rgb("#D63A57")  # rose
    C3       = hex_rgb("#1F9B8E")  # teal
    C4       = hex_rgb("#7DD4C9")  # mint (inner)

    C1_HEX, C2_HEX, C3_HEX, C4_HEX = "#E5A700", "#D63A57", "#1F9B8E", "#7DD4C9"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_rect(slide, 0, 0, 13.333, 7.5, BG)

    # Left accent stripe
    add_rect(slide, 0, 0, 0.25, 7.5, C3)

    # Title block
    add_text(slide, 0.7, 0.5, 11, 0.3, "TARGET AUDIENCES",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3)
    add_text(slide, 0.7, 0.85, 11, 0.8, args.title,
             font="Trebuchet MS", size=fit_title_pt(args.title, 11), bold=True, color=INK)
    type_label = {"sequel": "Sequel", "new_ip_with_fans": "IP-based",
                  "custom": "Original IP"}[args.type]
    add_text(slide, 0.7, 1.65, 11.3, 0.4,
             f"Potential buyer audience by tier · {args.genre} · {type_label}",
             font="Calibri", size=14, color=MUTED)

    # ---- Circles (left half) — filled, labeled tiers ----
    cx, cy = 3.75, 4.85
    d1, d2, d3, d4 = 5.3, 4.05, 2.75, 1.5
    add_circle(slide, cx, cy, d1, C1, line=BG, line_w_pt=2.5)
    add_circle(slide, cx, cy, d2, C2, line=BG, line_w_pt=2.5)
    add_circle(slide, cx, cy, d3, C3, line=BG, line_w_pt=2.5)
    add_circle(slide, cx, cy, d4, C4, line=BG, line_w_pt=2.5)

    inner_chart = ring2_chart_label(args)
    prev_chart, _ = innermost_label(args)

    fg1 = fg_for_bg(C1_HEX, "#FFFFFF", "#1A1A1A")
    fg2 = fg_for_bg(C2_HEX, "#FFFFFF", "#1A1A1A")
    fg3 = fg_for_bg(C3_HEX, "#FFFFFF", "#1A1A1A")
    fg4 = fg_for_bg(C4_HEX, "#FFFFFF", "#1A1A1A")

    r1, r2, r3, r4 = d1 / 2, d2 / 2, d3 / 2, d4 / 2

    # v9 pass: same arc-text-only treatment as the dark theme (see
    # render_ring_label) so both themes are agency-caliber and consistent.
    # Numeric counts live only in the side legend now.
    c4_label = cohort4_name(args)
    render_ring_label(slide, cx, cy, r1, r2, c4_label,
                       fg_color=fg1, base_pt=15, min_pt=9)

    c3_label = cohort3_name(args)
    render_ring_label(slide, cx, cy, r2, r3, c3_label,
                       fg_color=fg2, base_pt=13, min_pt=8)

    render_ring_label(slide, cx, cy, r3, r4, inner_chart,
                       fg_color=fg3, base_pt=11.5, min_pt=7)

    # ---- Inner circle (cohort 1) ----
    # v9: number removed (lives only in the side legend now) -- center
    # circle shows ONLY the cohort-name label, auto-shrunk/wrapped so it
    # NEVER breaks outside the circle bounds.
    render_center_label(slide, cx, cy, d4, prev_chart, fg_color=fg4)

    # ---- Legend (right half) — kept, per user's "keep side card" pick ----
    rx = 7.7
    add_text(slide, rx, 2.25, 5, 0.3, "POTENTIAL BUYERS BY TIER",
             font="Calibri", size=body_pt(L, 10), bold=True, color=C3)
    add_rect(slide, rx, 2.55, 5.0, 0.012, HAIR)

    _, inner_legend = innermost_label(args)
    inner_desc = "Players of the prior title" if args.inner == "prev" else \
                 "Direct followers of the developer" if args.inner == "dev" else \
                 (args.inner_definition or "")

    rows = [
        (C4, inner_legend,                inner_desc,                              fmt_num(args.prev)),
        (C3, ring2_legend_name(args),     ring2_legend_desc(args),                 fmt_num(args.ip_fans)),
        (C2, cohort3_name(args),          cohort3_desc(args),                      fmt_num(args.genre_fans)),
        (C1, cohort4_name(args),          cohort4_desc(args),                      fmt_num(args.breakout)),
    ]
    # Row heights are computed per-row (not a fixed rh=1.03) because
    # user-typed cohort names can be long enough to wrap to 2 lines --
    # a fixed single-line name box caused the name to visually collide
    # with the description text directly below it.
    name_w = 2.6  # width available before the right-aligned number column
    def _name_line_count(nm: str) -> int:
        return 2 if len(nm) > 22 else 1

    y = 2.75
    for i, (c, name, desc, num) in enumerate(rows):
        n_lines = _name_line_count(name)
        name_h = 0.32 if n_lines == 1 else 0.62
        desc_y_off = 0.36 if n_lines == 1 else 0.66
        row_h = desc_y_off + 0.38
        # Color swatch spans the full row height
        add_rect(slide, rx, y + 0.1, 0.2, row_h - 0.2, c)
        # Name (word_wrap on, so long names wrap onto a 2nd line automatically)
        add_text(slide, rx + 0.4, y + 0.02, name_w, name_h, name,
                 font="Trebuchet MS", size=14, bold=True, color=INK)
        # Description -- offset down when the name wrapped to 2 lines
        add_text(slide, rx + 0.4, y + desc_y_off, 3.5, 0.4, desc,
                 font="Calibri", size=body_pt(L, 9.5), color=MUTED)
        # Number -- widened box + slightly smaller font so large formatted
        # numbers (e.g. "10,000,000") don't wrap awkwardly.
        add_text(slide, rx + 3.1, y + 0.12, 1.9, 0.5, num,
                 font="Trebuchet MS", size=16, bold=True, color=INK,
                 align=PP_ALIGN.RIGHT)
        if i < 3:
            add_rect(slide, rx, y + row_h - 0.06, 5.0, 0.008, HAIR)
        y += row_h

    prs.save(out_path)


# ---------- PNG export ----------
def convert_to_png(pptx_path: str, png_path: str) -> str:
    out_dir = os.path.dirname(os.path.abspath(png_path))
    os.makedirs(out_dir, exist_ok=True)
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
            check=True, capture_output=True,
        )
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(tmp, base + ".pdf")
        prefix = os.path.join(tmp, "slide")
        subprocess.run(
            ["pdftoppm", "-png", "-r", "200", pdf_path, prefix],
            check=True, capture_output=True,
        )
        produced = None
        for fn in sorted(os.listdir(tmp)):
            if fn.startswith("slide-") and fn.endswith(".png"):
                produced = os.path.join(tmp, fn)
                break
        if not produced:
            raise RuntimeError("pdftoppm did not produce a PNG")
        with open(produced, "rb") as src, open(png_path, "wb") as dst:
            dst.write(src.read())
    return png_path


# ---------- CLI ----------
def parse_args():
    p = argparse.ArgumentParser()
    p.add_argument("--title", required=True)
    p.add_argument("--genre", required=True)
    p.add_argument(
        "--type", required=True,
        choices=["sequel", "new_ip_with_fans", "custom"],
        help="sequel | new_ip_with_fans | custom (custom = neither sequel nor IP-based; ring 2 is user-defined)",
    )
    p.add_argument(
        "--theme", required=True, choices=["dark", "light"],
        help="dark = V2 Modern Mono; light = V4 Bold Brand",
    )
    p.add_argument(
        "--inner", required=True, choices=["prev", "dev", "other"],
        help="Innermost circle cohort: prev (Prev Game Owners) | dev (Developer Fans) | other (custom)",
    )
    p.add_argument("--inner-name", default=None,
                   help="Custom innermost label (required when --inner other)")
    p.add_argument("--inner-definition", default=None,
                   help="Optional one-line definition shown in legend for the innermost cohort when --inner other")

    p.add_argument("--prev", type=int, required=True,
                   help="Innermost cohort size (prev owners / dev fans / custom inner)")
    p.add_argument("--ip-fans", type=int, required=True,
                   help="Ring 2 cohort size")
    p.add_argument("--genre-fans", type=int, required=True)
    p.add_argument("--breakout", type=int, required=True)

    p.add_argument("--ring2-name", default=None,
                   help="Custom cohort name for ring 2 (required when --type custom)")
    p.add_argument("--ring2-definition", default=None,
                   help="One-line definition shown in legend for ring 2 (required when --type custom)")

    p.add_argument("--cohort3-name", default=None,
                   help="User-typed name for cohort 3 (ring 2 from outside / 'genre fans' slot). "
                        "Falls back to 'Genre Fans' only if omitted -- the wizard always supplies this.")
    p.add_argument("--cohort3-definition", default=None,
                   help="Optional one-line definition for cohort 3 shown in the side legend.")
    p.add_argument("--cohort4-name", default=None,
                   help="User-typed name for cohort 4 (outer ring / 'breakout ceiling' slot). "
                        "Falls back to 'Breakout Ceiling' only if omitted -- the wizard always supplies this.")
    p.add_argument("--cohort4-definition", default=None,
                   help="Optional one-line definition for cohort 4 shown in the side legend.")

    p.add_argument("--out-dir", required=True)
    args = p.parse_args()

    if args.type == "custom" and (not args.ring2_name or not args.ring2_definition):
        p.error("--ring2-name and --ring2-definition are required when --type custom")
    if args.inner == "other" and not args.inner_name:
        p.error("--inner-name is required when --inner other")
    return args


def main():
    args = parse_args()
    slug = slugify(args.title)
    out_dir = os.path.abspath(args.out_dir)
    os.makedirs(out_dir, exist_ok=True)
    pptx_path = os.path.join(out_dir, f"{slug}_sizing_circle_{args.theme}.pptx")
    png_path  = os.path.join(out_dir, f"{slug}_sizing_circle_{args.theme}.png")

    if args.theme == "dark":
        render_dark(args, pptx_path)
    else:
        render_light(args, pptx_path)

    convert_to_png(pptx_path, png_path)
    print(f"PPTX: {pptx_path}")
    print(f"PNG:  {png_path}")


if __name__ == "__main__":
    main()
